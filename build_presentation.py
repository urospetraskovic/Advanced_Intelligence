"""
Generator PPTX prezentacije za predmet "Napredne tehnike računarske inteligencije".

Tema: LLM-driven pipeline za automatizovan Intelligent Tutoring System
Tim: Uroš Petrašković, Luka Šarić, Stefan Lazarević

Dizajn principi:
- Sve pozicije se izvode iz LAYOUT konstanti (responsive, simetrično, bez hardkodovanja).
- Helperi za naslov, sadržaj, tabele, citate, image placeholder-e, kod.
- Minimalistički: belina, jasna hijerarhija, akcentna boja samo za fokus.
- Image placeholder = svetlosivi okvir sa labelom (mesto da tim ubaci sliku).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy
from lxml import etree


# =============================================================================
# DIZAJN TOKENI
# =============================================================================

# Boje (akademski, minimalistički ton)
C_PRIMARY      = RGBColor(0x1B, 0x36, 0x5D)   # tamno teget (naslovi, akcenti)
C_ACCENT       = RGBColor(0xC2, 0x57, 0x00)   # topla bakarna (highlights, brojevi)
C_TEXT         = RGBColor(0x1F, 0x29, 0x37)   # skoro crna
C_MUTED        = RGBColor(0x4B, 0x55, 0x63)   # siva za sekundarni tekst (pojačan kontrast)
C_LIGHT_LINE   = RGBColor(0xE5, 0xE7, 0xEB)   # svetle linije i borderi
C_CARD_BG      = RGBColor(0xF5, 0xF6, 0xF8)   # svetla pozadina kartica
C_IMG_BG       = RGBColor(0xEE, 0xEF, 0xF2)   # placeholder pozadina
C_WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
C_TABLE_HEADER = RGBColor(0x1B, 0x36, 0x5D)
C_TABLE_ALT    = RGBColor(0xFA, 0xFB, 0xFC)

# Tipografija
FONT_FAMILY        = "Calibri"
FONT_FAMILY_MONO   = "Consolas"

FS_TITLE_BIG       = 40   # naslovna strana
FS_TITLE           = 30   # naslov slajda
FS_SECTION_LABEL   = 14   # "SEGMENT B" sitno gore
FS_SECTION_TITLE   = 44   # veliki naslov section dividera
FS_SECTION_SUB     = 18   # podnaslov section dividera
FS_SUBTITLE        = 18
FS_BODY            = 17
FS_BODY_SM         = 15
FS_CAPTION         = 11
FS_FOOTER          = 9
FS_CODE            = 13
FS_QUOTE           = 26

# Layout (16:9 widescreen)
SLIDE_W            = Inches(13.333)
SLIDE_H            = Inches(7.5)

# Spoljne margine
MARGIN_X           = Inches(0.6)
MARGIN_TOP         = Inches(0.45)
MARGIN_BOTTOM      = Inches(0.45)

# Header
HEADER_H           = Inches(0.95)   # visina naslovnog bloka
ACCENT_LINE_H      = Emu(20000)     # tanka akcentna linija ispod naslova
GAP_AFTER_HEADER   = Inches(0.18)

# Footer
FOOTER_H           = Inches(0.3)

# Sirina i visina kontent regiona
CONTENT_X          = MARGIN_X
CONTENT_W          = SLIDE_W - 2 * MARGIN_X
CONTENT_Y          = MARGIN_TOP + HEADER_H + GAP_AFTER_HEADER
CONTENT_H          = SLIDE_H - CONTENT_Y - MARGIN_BOTTOM - FOOTER_H

# Konstantne reference za prezentaciju
PROJECT_FOOTER     = "LLM-driven ITS pipeline - Napredne tehnike računarske inteligencije"


# =============================================================================
# LOW-LEVEL HELPER-I
# =============================================================================

def set_run(run, *, size=FS_BODY, bold=False, italic=False,
            color=C_TEXT, font=FONT_FAMILY):
    """Postavlja stil run-a (jedinog tekst-fragmenta)."""
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text_box(slide, x, y, w, h, *, text="", size=FS_BODY, bold=False,
                 italic=False, color=C_TEXT, font=FONT_FAMILY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, word_wrap=True):
    """Dodaje textbox sa jednim run-om."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = word_wrap
    tf.vertical_anchor = anchor

    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    set_run(run, size=size, bold=bold, italic=italic, color=color, font=font)
    return tb


def add_filled_rect(slide, x, y, w, h, *, fill=C_WHITE, line=None, line_w=0.75):
    """Pravougaonik sa fill bojom i opcionim border-om."""
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid()
    rect.fill.fore_color.rgb = fill
    if line is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = line
        rect.line.width = Pt(line_w)
    rect.shadow.inherit = False
    return rect


def add_horizontal_line(slide, x, y, w, *, color=C_PRIMARY, height=ACCENT_LINE_H):
    """Tanka horizontalna linija (kao akcent)."""
    return add_filled_rect(slide, x, y, w, height, fill=color, line=None)


def set_solid_bg(slide, color):
    """Postavlja punu pozadinsku boju slajda."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


# =============================================================================
# STANDARDNI BLOKOVI (HEADER, FOOTER, KARTICE)
# =============================================================================

def add_header(slide, *, title, subtitle=None, page_label=None):
    """
    Dosledan zaglavlje bloka:
      - mali "page label" (npr. SEGMENT B - ONTOLOGIJE)
      - veliki naslov
      - opcioni podnaslov
      - tanka akcentna linija
    """
    y = MARGIN_TOP

    if page_label:
        add_text_box(
            slide, MARGIN_X, y, CONTENT_W, Inches(0.22),
            text=page_label.upper(), size=FS_SECTION_LABEL,
            bold=True, color=C_ACCENT,
        )
        title_y = y + Inches(0.27)
    else:
        title_y = y

    add_text_box(
        slide, MARGIN_X, title_y, CONTENT_W, Inches(0.55),
        text=title, size=FS_TITLE, bold=True, color=C_PRIMARY,
    )

    if subtitle:
        add_text_box(
            slide, MARGIN_X, title_y + Inches(0.55), CONTENT_W, Inches(0.3),
            text=subtitle, size=FS_SUBTITLE, italic=True, color=C_MUTED,
        )


def add_footer(slide, *, page_number=None, total=None):
    """Diskretan footer: sitna labela dole levo (bez slide counter-a)."""
    y = SLIDE_H - FOOTER_H - Inches(0.05)
    add_text_box(
        slide, MARGIN_X, y, CONTENT_W, FOOTER_H,
        text=PROJECT_FOOTER, size=FS_FOOTER, color=C_MUTED,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def add_card(slide, x, y, w, h, *, title=None, body_lines=None,
             title_color=C_PRIMARY, accent_top=True):
    """Kartica sa svetlom pozadinom, opcionim naslovom i bulletima."""
    card = add_filled_rect(slide, x, y, w, h, fill=C_CARD_BG,
                           line=C_LIGHT_LINE, line_w=0.5)

    pad_x = Inches(0.22)
    pad_y = Inches(0.16)
    inner_x = x + pad_x
    inner_w = w - 2 * pad_x

    cur_y = y + pad_y

    if accent_top:
        add_horizontal_line(slide, x, y, w, color=C_PRIMARY,
                            height=Emu(25000))

    if title:
        add_text_box(
            slide, inner_x, cur_y, inner_w, Inches(0.3),
            text=title, size=FS_SUBTITLE, bold=True, color=title_color,
        )
        cur_y += Inches(0.36)

    if body_lines:
        remaining_h = (y + h) - cur_y - pad_y
        add_bullet_list(slide, inner_x, cur_y, inner_w, remaining_h,
                        body_lines, size=FS_BODY_SM)

    return card


def add_image_placeholder(slide, x, y, w, h, *, label="Slika"):
    """
    Vizuelni placeholder za buduće ubacivanje slike.
    Pravi svetlo-sivi okvir sa labelom u centru.
    """
    add_filled_rect(slide, x, y, w, h, fill=C_IMG_BG,
                    line=C_LIGHT_LINE, line_w=0.75)

    # ikonica simulirana karakterom + tekst
    icon_h = Inches(0.4)
    icon_tb = add_text_box(
        slide, x, y + h / 2 - icon_h, w, icon_h,
        text="[ Slika ]", size=20, bold=True, color=C_MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM,
    )

    add_text_box(
        slide, x, y + h / 2 + Inches(0.05), w, Inches(0.4),
        text=label, size=FS_CAPTION, italic=True, color=C_MUTED,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP,
    )


# =============================================================================
# BULLET LISTE I PARAGRAFI
# =============================================================================

def add_bullet_list(slide, x, y, w, h, items, *,
                    size=FS_BODY, color=C_TEXT, line_spacing=1.25,
                    space_after=4, bullet_char="•"):
    """
    Bullet lista. items može biti:
      - string: običan bullet
      - (text, "strong"): section header (uppercase, accent, bez crtice)
      - (text, "sub"): indented sub-bullet
    Prazan string "" pravi mali vertikalni razmak između grupa.
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)
    tf.word_wrap = True

    first = True
    for item in items:
        kind = None
        text = item
        if isinstance(item, tuple):
            text = item[0]
            kind = item[1] if len(item) > 1 else None

        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing

        if text == "":
            # Spacer između grupa
            p.space_after = Pt(6)
            p.space_before = Pt(0)
            run = p.add_run()
            run.text = " "
            set_run(run, size=4, color=color)
            continue

        if kind == "strong":
            # Section header: mala uppercase labela u akcentnoj boji
            p.space_before = Pt(4)
            p.space_after = Pt(3)
            run = p.add_run()
            run.text = text.upper()
            set_run(run, size=11, bold=True, color=C_ACCENT)
        elif kind == "sub":
            p.space_after = Pt(space_after)
            run = p.add_run()
            run.text = "      ◦  " + text
            set_run(run, size=size - 1, color=color)
        else:
            p.space_after = Pt(space_after)
            run = p.add_run()
            run.text = f"{bullet_char}  " + text
            set_run(run, size=size, color=color)

    return tb


def add_two_column_bullets(slide, x, y, w, h, left_title, left_items,
                            right_title, right_items, *, gap=Inches(0.35)):
    """Dve simetrične kolone sa naslovima i bulletima."""
    col_w = (w - gap) / 2

    # Levo
    add_text_box(slide, x, y, col_w, Inches(0.32),
                 text=left_title, size=FS_SUBTITLE, bold=True, color=C_PRIMARY)
    add_horizontal_line(slide, x, y + Inches(0.36), Inches(0.6),
                        color=C_ACCENT, height=Emu(20000))
    add_bullet_list(slide, x, y + Inches(0.5), col_w, h - Inches(0.5),
                    left_items, size=FS_BODY)

    # Desno
    rx = x + col_w + gap
    add_text_box(slide, rx, y, col_w, Inches(0.32),
                 text=right_title, size=FS_SUBTITLE, bold=True, color=C_PRIMARY)
    add_horizontal_line(slide, rx, y + Inches(0.36), Inches(0.6),
                        color=C_ACCENT, height=Emu(20000))
    add_bullet_list(slide, rx, y + Inches(0.5), col_w, h - Inches(0.5),
                    right_items, size=FS_BODY)


# =============================================================================
# TABELE
# =============================================================================

def add_styled_table(slide, x, y, w, h, headers, rows, *,
                     col_widths=None, header_color=C_TABLE_HEADER,
                     header_text=C_WHITE, alt=True):
    """
    Stilizovana tabela: header trake u akcentnoj boji, alternativne redove
    u svetloj boji, tanke linije.
    """
    n_cols = len(headers)
    n_rows = len(rows) + 1
    table_shape = slide.shapes.add_table(n_rows, n_cols, x, y, w, h)
    table = table_shape.table

    # Sirine kolona
    if col_widths:
        total = sum(col_widths)
        for i, ratio in enumerate(col_widths):
            table.columns[i].width = int(w * ratio / total)

    # Header
    for i, htext in enumerate(headers):
        cell = table.cell(0, i)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_color
        cell.margin_left = cell.margin_right = Inches(0.1)
        cell.margin_top = cell.margin_bottom = Inches(0.06)
        tf = cell.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = htext
        set_run(run, size=FS_BODY_SM, bold=True, color=header_text)

    # Rows
    for r, row in enumerate(rows, start=1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_TABLE_ALT if (alt and r % 2 == 0) else C_WHITE
            cell.margin_left = cell.margin_right = Inches(0.1)
            cell.margin_top = cell.margin_bottom = Inches(0.05)
            tf = cell.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT
            run = p.add_run()
            run.text = str(val)
            set_run(run, size=FS_BODY_SM, color=C_TEXT)

    return table


# =============================================================================
# CODE BLOK
# =============================================================================

def add_code_block(slide, x, y, w, h, code_text, *, caption=None):
    """Tamniji blok sa monospace tekstom (za code/SPARQL/JSON snippete)."""
    bg = add_filled_rect(slide, x, y, w, h, fill=RGBColor(0x1F, 0x29, 0x37))
    pad = Inches(0.18)
    tb = slide.shapes.add_textbox(x + pad, y + pad, w - 2 * pad, h - 2 * pad)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0)
    tf.margin_top = tf.margin_bottom = Inches(0)

    for i, line in enumerate(code_text.splitlines()):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        run.text = line if line else " "
        set_run(run, size=FS_CODE, color=RGBColor(0xE5, 0xE7, 0xEB),
                font=FONT_FAMILY_MONO)

    if caption:
        add_text_box(slide, x, y + h + Inches(0.05), w, Inches(0.25),
                     text=caption, size=FS_CAPTION, italic=True, color=C_MUTED)


# =============================================================================
# SHEME SLAJDOVA (TEMPLATES)
# =============================================================================

def new_slide(prs):
    """Prazan layout (bez placeholdera) i bela pozadina."""
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    set_solid_bg(slide, C_WHITE)
    return slide


def make_title_slide(prs, *, title, subtitle, team, course, date_str):
    slide = new_slide(prs)
    set_solid_bg(slide, C_WHITE)

    # Tanka leva traka (akcent)
    add_filled_rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H,
                    fill=C_PRIMARY)

    # Course label (gore)
    add_text_box(
        slide, Inches(0.85), Inches(0.7), SLIDE_W - Inches(1.6), Inches(0.3),
        text=course.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
    )

    # Naslov
    add_text_box(
        slide, Inches(0.85), Inches(1.5), SLIDE_W - Inches(1.6), Inches(2.0),
        text=title, size=FS_TITLE_BIG, bold=True, color=C_PRIMARY,
    )

    # Akcentna linija
    add_horizontal_line(slide, Inches(0.85), Inches(3.7), Inches(1.2),
                        color=C_ACCENT, height=Emu(30000))

    # Podnaslov
    add_text_box(
        slide, Inches(0.85), Inches(3.9), SLIDE_W - Inches(1.6), Inches(1.2),
        text=subtitle, size=20, italic=True, color=C_MUTED,
    )

    # Tim
    add_text_box(
        slide, Inches(0.85), Inches(6.0), SLIDE_W - Inches(1.6), Inches(0.4),
        text="Tim", size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
    )
    add_text_box(
        slide, Inches(0.85), Inches(6.3), SLIDE_W - Inches(1.6), Inches(0.4),
        text=team, size=FS_BODY, color=C_TEXT,
    )

    # Datum
    add_text_box(
        slide, Inches(0.85), Inches(6.8), SLIDE_W - Inches(1.6), Inches(0.3),
        text=date_str, size=FS_CAPTION, color=C_MUTED,
    )

    return slide


def make_section_divider(prs, *, segment_label, title, subtitle, time_str):
    slide = new_slide(prs)
    set_solid_bg(slide, C_PRIMARY)

    # Akcentna pruga
    add_filled_rect(slide, Inches(0), SLIDE_H - Inches(0.4), SLIDE_W, Inches(0.4),
                    fill=C_ACCENT)

    add_text_box(
        slide, Inches(0.85), Inches(2.4), SLIDE_W - Inches(1.7), Inches(0.4),
        text=segment_label.upper(), size=18, bold=True,
        color=RGBColor(0xF5, 0xC0, 0x60),
    )

    add_text_box(
        slide, Inches(0.85), Inches(2.85), SLIDE_W - Inches(1.7), Inches(1.6),
        text=title, size=FS_SECTION_TITLE, bold=True, color=C_WHITE,
    )

    add_horizontal_line(slide, Inches(0.85), Inches(4.55), Inches(1.5),
                        color=C_ACCENT, height=Emu(40000))

    add_text_box(
        slide, Inches(0.85), Inches(4.75), SLIDE_W - Inches(1.7), Inches(1.0),
        text=subtitle, size=FS_SECTION_SUB, italic=True,
        color=RGBColor(0xD1, 0xD5, 0xDB),
    )

    add_text_box(
        slide, Inches(0.85), SLIDE_H - Inches(1.0), SLIDE_W - Inches(1.7), Inches(0.4),
        text=time_str, size=FS_CAPTION, color=RGBColor(0xC4, 0xC9, 0xD4),
    )

    return slide


def make_content_slide(prs, *, page_label, title, subtitle=None,
                       body_func=None, page_no=None, total_pages=None):
    slide = new_slide(prs)
    add_header(slide, title=title, subtitle=subtitle, page_label=page_label)
    if body_func:
        body_func(slide)
    add_footer(slide, page_number=page_no, total=total_pages)
    return slide


def make_quote_slide(prs, *, page_label, quote, attribution=None,
                     page_no=None, total_pages=None):
    slide = new_slide(prs)

    # Tanak akcent gore
    add_horizontal_line(slide, Inches(0), Inches(0), SLIDE_W,
                        color=C_PRIMARY, height=Emu(40000))

    # Sitan label
    add_text_box(
        slide, MARGIN_X, Inches(0.9), CONTENT_W, Inches(0.3),
        text=page_label.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
    )

    # Veliki citat
    add_text_box(
        slide, Inches(1.5), Inches(2.4), SLIDE_W - Inches(3.0), Inches(3.0),
        text=f"“{quote}”", size=FS_QUOTE, italic=True, color=C_PRIMARY,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
    )

    if attribution:
        add_text_box(
            slide, Inches(1.5), Inches(5.5), SLIDE_W - Inches(3.0), Inches(0.4),
            text=f"- {attribution}", size=FS_BODY, color=C_MUTED,
            align=PP_ALIGN.CENTER,
        )

    add_footer(slide, page_number=page_no, total=total_pages)
    return slide


def make_thank_you_slide(prs, *, team):
    slide = new_slide(prs)
    set_solid_bg(slide, C_PRIMARY)

    add_text_box(
        slide, Inches(0.85), Inches(2.8), SLIDE_W - Inches(1.7), Inches(1.4),
        text="Hvala na pažnji.", size=54, bold=True, color=C_WHITE,
        align=PP_ALIGN.CENTER,
    )

    add_horizontal_line(slide, SLIDE_W / 2 - Inches(0.6), Inches(4.2), Inches(1.2),
                        color=C_ACCENT, height=Emu(40000))

    add_text_box(
        slide, Inches(0.85), Inches(4.5), SLIDE_W - Inches(1.7), Inches(0.6),
        text="Pitanja i diskusija", size=24, italic=True,
        color=RGBColor(0xE5, 0xE7, 0xEB), align=PP_ALIGN.CENTER,
    )

    add_text_box(
        slide, Inches(0.85), Inches(6.0), SLIDE_W - Inches(1.7), Inches(0.4),
        text=team, size=FS_BODY, color=RGBColor(0xC4, 0xC9, 0xD4),
        align=PP_ALIGN.CENTER,
    )

    add_filled_rect(slide, Inches(0), SLIDE_H - Inches(0.3), SLIDE_W, Inches(0.3),
                    fill=C_ACCENT)
    return slide


# =============================================================================
# SADRZAJ - SLAJDOVI
# =============================================================================
# Svaki body_* helper prima 'slide' i koristi CONTENT_X / CONTENT_Y /
# CONTENT_W / CONTENT_H za simetričan raspored.

def body_agenda(slide):
    items = [
        ("Uvod i motivacija",            "10 min", "Šta i zašto - od ishoda učenja do ITS-a"),
        ("LLM prompting i taksonomije",  "25 min", "Bloom, SOLO, PS4 prompting, hijerarhija zadataka"),
        ("Ontologije i formalni grounding","25 min", "OWL, RDF, SPARQL i anti-halucinacioni mehanizmi"),
        ("Evaluacija i ograničenja",     "25 min", "Kvalitet po nivoima, robustnost, human-in-the-loop"),
        ("Sinteza, budući rad, zaključak","10-15 min","Integrisana arhitektura i naredni koraci"),
    ]

    n = len(items)
    gap = Inches(0.12)
    row_h = (CONTENT_H - (n - 1) * gap) / n

    for i, (title, time, sub) in enumerate(items):
        y = CONTENT_Y + i * (row_h + gap)

        # leva traka
        add_filled_rect(slide, CONTENT_X, y, Inches(0.12), row_h,
                        fill=C_ACCENT)

        # broj
        add_text_box(
            slide, CONTENT_X + Inches(0.25), y, Inches(0.7), row_h,
            text=f"0{i}.", size=24, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )

        # naslov
        add_text_box(
            slide, CONTENT_X + Inches(1.05), y + Inches(0.05),
            CONTENT_W - Inches(3.0), Inches(0.45),
            text=title, size=20, bold=True, color=C_PRIMARY,
        )

        # podnaslov
        add_text_box(
            slide, CONTENT_X + Inches(1.05), y + Inches(0.5),
            CONTENT_W - Inches(3.0), row_h - Inches(0.5),
            text=sub, size=FS_BODY_SM, color=C_MUTED,
        )

        # vreme desno
        add_text_box(
            slide, CONTENT_X + CONTENT_W - Inches(1.9), y,
            Inches(1.9), row_h,
            text=time, size=FS_BODY, bold=True, color=C_ACCENT,
            align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE,
        )


def body_team(slide):
    members = [
        ("Uroš Petrašković",
         "SOLO Quiz Generator",
         "Generisanje pitanja po SOLO taksonomiji uz ontološki grounding i source-line citate."),
        ("Luka Šarić",
         "Course Structure Generator",
         "Automatska struktura kursa po Bloomovoj taksonomiji - moduli, koncepti, ishodi."),
        ("Stefan Lazarević",
         "Computer Use Video Tutorials",
         "Video demonstracije praktičkih zadataka uz OWL validaciju izvršnog plana."),
    ]

    gap = Inches(0.3)
    col_w = (CONTENT_W - 2 * gap) / 3

    for i, (name, role, desc) in enumerate(members):
        x = CONTENT_X + i * (col_w + gap)

        # kartica
        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.3),
                        fill=C_CARD_BG, line=C_LIGHT_LINE, line_w=0.75)

        # akcent na vrhu kartice
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.25)
        inner_x = x + pad
        inner_w = col_w - 2 * pad
        cur_y = CONTENT_Y + Inches(0.3)

        # avatar krug placeholder
        avatar_size = Inches(1.0)
        avatar_x = x + (col_w - avatar_size) / 2
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, avatar_x, cur_y,
                                        avatar_size, avatar_size)
        circle.fill.solid()
        circle.fill.fore_color.rgb = C_PRIMARY
        circle.line.fill.background()
        # inicijal
        ini = name.split()[0][0] + name.split()[1][0]
        tb = circle.text_frame
        tb.margin_left = tb.margin_right = Inches(0)
        tb.margin_top = tb.margin_bottom = Inches(0)
        p = tb.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = ini
        set_run(run, size=28, bold=True, color=C_WHITE)

        cur_y += avatar_size + Inches(0.2)

        add_text_box(
            slide, inner_x, cur_y, inner_w, Inches(0.4),
            text=name, size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
            align=PP_ALIGN.CENTER,
        )
        cur_y += Inches(0.4)

        add_text_box(
            slide, inner_x, cur_y, inner_w, Inches(0.35),
            text=role, size=FS_BODY_SM, italic=True, color=C_ACCENT,
            align=PP_ALIGN.CENTER,
        )
        cur_y += Inches(0.5)

        add_text_box(
            slide, inner_x, cur_y, inner_w, CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.3),
            text=desc, size=FS_BODY_SM, color=C_TEXT,
            align=PP_ALIGN.CENTER,
        )


def body_predmet_context(slide):
    # Leva polovina - tekst
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y, half_w, CONTENT_H,
        [
            ("Predmet: Napredne tehnike računarske inteligencije", "strong"),
            "Master studije, FTN Univerzitet u Novom Sadu",
            "Predavači: dr Prokić i dr Kovačević",
            "",
            ("Tema u preseku dve oblasti", "strong"),
            "AI in (Computing) Education",
            "Intelligent Tutoring Systems",
            "",
            ("Format polaganja", "strong"),
            "Snimljena prezentacija 1h - 1h30min",
            "Demonstracija integrisanog ITS pipeline-a",
        ],
        size=FS_BODY,
    )

    # Desna polovina - vizuelni placeholder
    right_x = CONTENT_X + half_w + Inches(0.4)
    add_image_placeholder(
        slide, right_x, CONTENT_Y, half_w, CONTENT_H,
        label="Logo FTN-a / vizual predmeta",
    )


def body_problem_stats(slide):
    """3 brojke + sažet opis problema."""
    # Gornji red - 3 velike brojke
    stats_h = Inches(2.1)
    gap = Inches(0.3)
    col_w = (CONTENT_W - 2 * gap) / 3

    stats = [
        ("~20%",
         "radnog vremena nastavnika",
         "troši se na pripremu i administraciju (OECD TALIS 2018, prosek)"),
        ("7-12 h",
         "nedeljno na pripremu",
         "van časa, samo na materijale (NCES 2019 Teacher Time Study)"),
        ("22 min",
         "po jednom MCQ pitanju",
         "sa distraktorima i mapiranjem na ishod (Walker et al., 2020)"),
    ]

    for i, (num, lbl, sub) in enumerate(stats):
        x = CONTENT_X + i * (col_w + gap)
        add_filled_rect(slide, x, CONTENT_Y, col_w, stats_h,
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, CONTENT_Y, Inches(0.08), stats_h,
                        fill=C_ACCENT)

        add_text_box(
            slide, x + Inches(0.3), CONTENT_Y + Inches(0.2),
            col_w - Inches(0.5), Inches(0.9),
            text=num, size=38, bold=True, color=C_ACCENT,
        )
        add_text_box(
            slide, x + Inches(0.3), CONTENT_Y + Inches(1.05),
            col_w - Inches(0.5), Inches(0.35),
            text=lbl, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, x + Inches(0.3), CONTENT_Y + Inches(1.4),
            col_w - Inches(0.5), Inches(0.7),
            text=sub, size=FS_BODY_SM, color=C_MUTED,
        )

    # Donji blok - poruka
    y2 = CONTENT_Y + stats_h + Inches(0.4)
    add_text_box(
        slide, CONTENT_X, y2, CONTENT_W, Inches(0.4),
        text="Šta nastavnik radi ručno?", size=FS_SUBTITLE, bold=True,
        color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, y2 + Inches(0.45), CONTENT_W,
        CONTENT_H - (y2 + Inches(0.45) - CONTENT_Y),
        [
            "Planira strukturu kursa - module, koncepte, ishode učenja",
            "Piše objašnjenja, primere, definicije - usaglašava ih sa nivoom studenata",
            "Priprema praktične demonstracije, često u formi video tutorijala",
            "Sastavlja pitanja na različitim kognitivnim nivoima - i tu nastaje uska grla",
        ],
        size=FS_BODY,
    )


def body_llm_kao_odgovor(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="LLM je otvorio vrata", size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.45), half_w, CONTENT_H - Inches(0.5),
        [
            "Generisanje sadržaja iz prirodno-jezičkog uputstva",
            "Razlaganje kompleksnih zadataka",
            "Brza iteracija i prilagođavanje materijala",
            "Multimodalnost: tekst, kod, vizija",
        ],
        size=FS_BODY,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="...ali otvorio i nove probleme", size=FS_SUBTITLE, bold=True,
        color=C_ACCENT,
    )
    add_bullet_list(
        slide, right_x, CONTENT_Y + Inches(0.45), half_w, CONTENT_H - Inches(0.5),
        [
            "Halucinacije - izmišljene činjenice",
            "Generičnost izlaza, malo veze sa konkretnim materijalom",
            "Pitanja na pogrešnom kognitivnom nivou",
            "Nedostatak formalne strukture i ponovljivosti",
            "Teško proveriti gde se šta oslonilo na izvor",
        ],
        size=FS_BODY,
    )


def body_why_now(slide):
    """Tri istorijska / empirijska argumenta za 'zašto sad'."""
    cards = [
        ("Bloom (1984)",
         "'2 Sigma Problem'",
         "1-na-1 tutoring podiže prosečnog studenta na ~98. percentil. "
         "ITS poslednjih 40 godina pokušava da se približi ovom efektu."),
        ("Khan Academy (2023+)",
         "Khanmigo u školama",
         "Prvi mainstream LLM tutor u nastavi. Već u 100+ školskih distrikta "
         "u SAD. AI tutori više nisu istraživanje - produkcija su."),
        ("Kestin et al. (Harvard, 2024)",
         "AI tutor > active learning",
         "U kontrolisanoj studiji iz fizike, GPT-4 tutor je nadmašio "
         "aktivno-učeničku nastavu, i to za pola vremena. Empirijski signal "
         "da je tehnologija sazrela."),
    ]
    gap = Inches(0.3)
    col_w = (CONTENT_W - 2 * gap) / 3

    for i, (src, title, body) in enumerate(cards):
        x = CONTENT_X + i * (col_w + gap)
        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.4),
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.25)
        cur_y = CONTENT_Y + Inches(0.3)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.3),
            text=src.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        )
        cur_y += Inches(0.35)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.7),
            text=title, size=20, bold=True, color=C_PRIMARY,
        )
        cur_y += Inches(0.7)
        add_horizontal_line(slide, x + pad, cur_y, Inches(0.5),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.2)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad,
            CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.4),
            text=body, size=FS_BODY_SM, color=C_TEXT,
        )

    # Donji red - sažeta poruka
    y2 = CONTENT_Y + CONTENT_H - Inches(0.3)
    add_text_box(
        slide, CONTENT_X, y2, CONTENT_W, Inches(0.3),
        text="Pitanje više nije 'da li LLM može nešto' - pitanje je 'kako ga koristimo dobro'.",
        size=FS_BODY_SM, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER,
    )


def body_its_landscape(slide):
    """Brzi pregled ITS pejzaza - gde se mi uklapamo."""
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="ITS landscape - gde se mi uklapamo",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    items = [
        ("Carnegie Learning / ALEKS",
         "Klasični cognitive tutor",
         "~30 godina, rule-based, jaka pedagoška teorija (ACT-R), bez LLM-a.",
         C_MUTED),
        ("Khan Academy Khanmigo",
         "LLM-driven, 2023+",
         "GPT-4 zasnovan tutor; čuvar konteksta + nastavnički dashboard.",
         C_PRIMARY),
        ("Duolingo Max",
         "LLM-driven jezici, 2023+",
         "'Role play' i 'Explain my answer' funkcije; placeno premium tier.",
         C_PRIMARY),
        ("Vesin et al. - ProTuS (FTN/UNS)",
         "Ontology-based Java tutor",
         "Domaci rad iz naseg fakulteta; ontologija + adaptivne preporuke.",
         C_PRIMARY),
        ("Nas predlog",
         "LLM + Ontologija + Taksonomija",
         "Generisanje materijala (struktura + pitanja + video), ne samo serviranje.",
         C_ACCENT),
    ]

    n = len(items)
    gap = Inches(0.1)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for name, kind, desc, color in items:
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.12), box_h, fill=color)
        add_text_box(
            slide, CONTENT_X + Inches(0.3), y, Inches(3.7), box_h,
            text=name, size=FS_BODY, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, CONTENT_X + Inches(4.1), y, Inches(2.6), box_h,
            text=kind, size=FS_BODY_SM, italic=True, color=color,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, CONTENT_X + Inches(6.8), y, CONTENT_W - Inches(7.0), box_h,
            text=desc, size=FS_BODY_SM, color=C_TEXT,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        y += box_h + gap


def body_vizija_pipeline(slide):
    # Pipeline kao 5 koraka
    steps = [
        ("Ulaz",          "Ishodi učenja\nPDF materijal\nNL instrukcije"),
        ("Struktura",     "Bloom-driven\nmoduli i ishodi"),
        ("Ontologija",    "OWL/RDF\nkonceptualne relacije"),
        ("Generisanje",   "SOLO pitanja\nvideo demoi"),
        ("Izlaz",         "Strukturisan\nITS pipeline"),
    ]
    n = len(steps)
    gap = Inches(0.18)
    box_h = Inches(2.4)
    box_w = (CONTENT_W - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.2)

    for i, (title, sub) in enumerate(steps):
        x = CONTENT_X + i * (box_w + gap)
        # kartica
        add_filled_rect(slide, x, y, box_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, y, box_w, Inches(0.1), fill=C_PRIMARY)

        # broj
        add_text_box(
            slide, x, y + Inches(0.2), box_w, Inches(0.4),
            text=f"0{i+1}", size=22, bold=True, color=C_ACCENT,
            align=PP_ALIGN.CENTER,
        )
        # naslov
        add_text_box(
            slide, x, y + Inches(0.7), box_w, Inches(0.4),
            text=title, size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
            align=PP_ALIGN.CENTER,
        )
        # podsadrzaj
        add_text_box(
            slide, x + Inches(0.1), y + Inches(1.2), box_w - Inches(0.2),
            box_h - Inches(1.3),
            text=sub, size=FS_BODY_SM, color=C_TEXT,
            align=PP_ALIGN.CENTER,
        )

        # strelica (osim za poslednji)
        if i < n - 1:
            arrow_x = x + box_w + Inches(0.02)
            arrow_y = y + box_h / 2 - Inches(0.1)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y,
                gap - Inches(0.04), Inches(0.2),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C_ACCENT
            arrow.line.fill.background()

    # Ispod - ključna poruka
    y2 = y + box_h + Inches(0.3)
    add_filled_rect(slide, CONTENT_X, y2, CONTENT_W, CONTENT_H - (y2 - CONTENT_Y) - Inches(0.1),
                    fill=C_CARD_BG, line=None)
    add_filled_rect(slide, CONTENT_X, y2, Inches(0.08),
                    CONTENT_H - (y2 - CONTENT_Y) - Inches(0.1), fill=C_ACCENT)
    add_text_box(
        slide, CONTENT_X + Inches(0.3), y2 + Inches(0.15),
        CONTENT_W - Inches(0.5), Inches(0.4),
        text="Vizija", size=FS_BODY_SM, bold=True, color=C_ACCENT,
    )
    add_text_box(
        slide, CONTENT_X + Inches(0.3), y2 + Inches(0.5),
        CONTENT_W - Inches(0.5), Inches(0.8),
        text="Od ishoda učenja, automatski, do snimljene video demonstracije i SOLO kviza - bez rucne intervencije.",
        size=FS_BODY, color=C_TEXT,
    )


def body_tri_komponente(slide):
    cols = [
        ("Luka",
         "Struktura kursa",
         "Bloomova taksonomija",
         ["Hijerarhija: kurs -> modul -> koncept -> ishod -> aktivnost",
          "JSON serijalizacija sa rekurzivnom strukturom",
          "Bloomovi nivoi mapirani na tipove aktivnosti",
          "State-of-generation kontekst za koherentnost"]),
        ("Stefan",
         "Video demonstracije",
         "OWL ontologija + Computer Use",
         ["NL instrukcija -> Task / Step / Action",
          "OWL validacija plana pre izvršenja",
          "SPARQL upiti vode izvršavanje",
          "Vision (Qwen 2.5 VL) za UI detekciju",
          "FFmpeg snima MP4 tutorijal"]),
        ("Uroš",
         "SOLO pitanja",
         "Ontology-grounded MCQ",
         ["4 SOLO nivoa: U / M / R / EA",
          "PS4 prompt template",
          "source_line citat iz PDF-a",
          "Ontology anchor iz ConceptRelationship",
          "Two-pass generisanje za EA"]),
    ]

    gap = Inches(0.25)
    col_w = (CONTENT_W - 2 * gap) / 3

    for i, (member, title, framework, points) in enumerate(cols):
        x = CONTENT_X + i * (col_w + gap)

        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.1),
                        fill=C_WHITE, line=C_LIGHT_LINE)
        # gornja akcent traka
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.2)
        cur_y = CONTENT_Y + Inches(0.25)

        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.3),
            text=member.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        )
        cur_y += Inches(0.32)

        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.45),
            text=title, size=22, bold=True, color=C_PRIMARY,
        )
        cur_y += Inches(0.5)

        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.3),
            text=framework, size=FS_BODY_SM, italic=True, color=C_MUTED,
        )
        cur_y += Inches(0.4)

        add_horizontal_line(slide, x + pad, cur_y, Inches(0.5),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.18)

        add_bullet_list(
            slide, x + pad, cur_y, col_w - 2 * pad,
            CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.3),
            points, size=FS_BODY_SM,
        )


def body_kljucna_poruka_uvod(slide):
    # Centriran citat sa "formula" stilom
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), CONTENT_W, Inches(0.5),
        text="Ključna poruka prezentacije", size=FS_SECTION_LABEL, bold=True,
        color=C_ACCENT, align=PP_ALIGN.CENTER,
    )

    # Formula u kartici
    box_w = CONTENT_W * 0.85
    box_x = CONTENT_X + (CONTENT_W - box_w) / 2
    box_y = CONTENT_Y + Inches(1.2)
    box_h = Inches(2.2)

    add_filled_rect(slide, box_x, box_y, box_w, box_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_horizontal_line(slide, box_x, box_y, box_w, color=C_PRIMARY,
                        height=Emu(30000))

    # Formula liniju po liniju
    add_text_box(
        slide, box_x, box_y + Inches(0.4), box_w, Inches(0.6),
        text="LLM  +  ontologija  +  pedagoška taksonomija",
        size=32, bold=True, color=C_PRIMARY,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide, box_x, box_y + Inches(1.1), box_w, Inches(0.45),
        text="=", size=28, bold=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    add_text_box(
        slide, box_x, box_y + Inches(1.55), box_w, Inches(0.5),
        text="pouzdano, pedagoški kalibrisano ITS okruženje",
        size=22, italic=True, color=C_TEXT,
        align=PP_ALIGN.CENTER,
    )

    # Sub
    add_text_box(
        slide, CONTENT_X, box_y + box_h + Inches(0.25), CONTENT_W, Inches(0.5),
        text="više od zbira svojih delova",
        size=FS_SUBTITLE, italic=True, color=C_MUTED,
        align=PP_ALIGN.CENTER,
    )


# --- SEGMENT A helperi -------------------------------------------------------

def body_bloom_pregled(slide):
    levels = [
        ("Pamtiti",      "Cinjenice, definicije, terminologija"),
        ("Razumeti",     "Objašnjenje, parafraziranje"),
        ("Primeniti",    "Resavanje analognih problema"),
        ("Analizirati",  "Razlaganje, poređenje"),
        ("Evaluirati",   "Procena, kritičko mišljenje"),
        ("Kreirati",     "Sinteza, novi rad"),
    ]
    half_w = CONTENT_W * 0.55
    n = len(levels)
    gap = Inches(0.08)
    row_h = (CONTENT_H - (n - 1) * gap) / n

    for i, (name, desc) in enumerate(levels):
        y = CONTENT_Y + i * (row_h + gap)
        add_filled_rect(slide, CONTENT_X, y, half_w, row_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        # broj badge
        add_filled_rect(slide, CONTENT_X, y, Inches(0.55), row_h, fill=C_PRIMARY)
        add_text_box(
            slide, CONTENT_X, y, Inches(0.55), row_h,
            text=f"L{i+1}", size=FS_BODY, bold=True, color=C_WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.7), y, half_w - Inches(0.8), Inches(0.35),
            text=name, size=FS_BODY, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.7), y + Inches(0.3),
            half_w - Inches(0.8), row_h - Inches(0.3),
            text=desc, size=FS_BODY_SM, color=C_MUTED,
        )

    # Desna polovina - placeholder za piramidu
    right_x = CONTENT_X + half_w + Inches(0.3)
    right_w = CONTENT_W - half_w - Inches(0.3)
    add_image_placeholder(slide, right_x, CONTENT_Y, right_w, CONTENT_H,
                          label="Bloomova piramida")


def body_solo_pregled(slide):
    levels = [
        ("Unistructural",
         "Jedna izolovana činjenica",
         "Setiti se jednog elementa"),
        ("Multistructural",
         "Više nepovezanih činjenica",
         "Prepoznati skup elemenata"),
        ("Relational",
         "Veza između koncepata",
         "Razumeti strukturu i odnose"),
        ("Extended Abstract",
         "Generalizacija i transfer",
         "Primeniti na novi kontekst"),
    ]
    half_w = CONTENT_W * 0.55
    n = len(levels)
    gap = Inches(0.12)
    row_h = (CONTENT_H - (n - 1) * gap) / n

    for i, (name, sub, desc) in enumerate(levels):
        y = CONTENT_Y + i * (row_h + gap)
        add_filled_rect(slide, CONTENT_X, y, half_w, row_h,
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.1), row_h, fill=C_ACCENT)
        add_text_box(
            slide, CONTENT_X + Inches(0.25), y + Inches(0.1),
            half_w - Inches(0.4), Inches(0.32),
            text=name, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.25), y + Inches(0.42),
            half_w - Inches(0.4), Inches(0.3),
            text=sub, size=FS_BODY_SM, italic=True, color=C_ACCENT,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.25), y + Inches(0.72),
            half_w - Inches(0.4), row_h - Inches(0.75),
            text=desc, size=FS_BODY_SM, color=C_MUTED,
        )

    right_x = CONTENT_X + half_w + Inches(0.3)
    right_w = CONTENT_W - half_w - Inches(0.3)
    add_image_placeholder(slide, right_x, CONTENT_Y, right_w, CONTENT_H - Inches(0.5),
                          label="SOLO taksonomija - vizuelni prikaz (Biggs & Tang)")
    # Citatna napomena ispod
    add_text_box(
        slide, right_x, CONTENT_Y + CONTENT_H - Inches(0.45), right_w, Inches(0.4),
        text="Biggs & Collis (1982). Lister et al. (2006) - empirijski "
             "pokazali da početnici programiranja ozbiljno zaostaju na "
             "Relational i Extended Abstract zadacima.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )


def body_bloom_vs_solo(slide):
    headers = ["Aspekt", "Bloom (revidirana)", "SOLO"]
    rows = [
        ["Tip hijerarhije",      "Tipovi misaonih operacija", "Strukturalna složenost odgovora"],
        ["Broj nivoa",           "6 nivoa", "5 nivoa (uključujući prestructural)"],
        ["Osnovna jedinica",     "Glagol akcije (znati, primeniti, ...)", "Broj i organizacija elemenata u odgovoru"],
        ["Tipična upotreba",     "Planiranje aktivnosti i ishoda učenja", "Evaluacija dubine učenja"],
        ["Mi koristimo za",      "Strukturu kursa (Luka)", "Generisanje pitanja (Uroš)"],
        ["Glavna prednost",      "Direktna veza sa nastavnim aktivnostima", "Direktna veza sa kognitivnom dubinom"],
    ]
    add_styled_table(slide, CONTENT_X, CONTENT_Y, CONTENT_W,
                     CONTENT_H - Inches(1.4), headers, rows,
                     col_widths=[1, 2, 2])

    # Donja poruka
    y2 = CONTENT_Y + CONTENT_H - Inches(1.2)
    add_filled_rect(slide, CONTENT_X, y2, CONTENT_W, Inches(1.0),
                    fill=C_CARD_BG, line=None)
    add_filled_rect(slide, CONTENT_X, y2, Inches(0.08), Inches(1.0), fill=C_ACCENT)
    add_text_box(
        slide, CONTENT_X + Inches(0.25), y2 + Inches(0.1),
        CONTENT_W - Inches(0.4), Inches(0.4),
        text="Zašto obe?", size=FS_BODY, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X + Inches(0.25), y2 + Inches(0.45),
        CONTENT_W - Inches(0.4), Inches(0.6),
        text="Bloom upravlja STA generišemo (tipovi aktivnosti), SOLO upravlja KOLIKO DUBOKO se proverava (kognitivni nivoi pitanja).",
        size=FS_BODY_SM, color=C_TEXT,
    )


def body_luka_struktura(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    # Levo - opis pristupa
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Hijerarhijska dekompozicija kursa",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Ulaz", "strong"),
            "Tema kursa i opšti opis",
            "Lista ciljeva (ishoda učenja)",
            "",
            ("LLM prolaz po nivou", "strong"),
            "Generiše module iz teme",
            "Generiše koncepte iz modula",
            "Generiše ishode iz koncepata",
            "Generiše aktivnosti iz ishoda",
            "",
            ("Bloomovo mapiranje", "strong"),
            "Svaka aktivnost dobija Bloomov nivo",
            "Aktivnosti se biraju kako bi pokrile sve nivoe",
        ],
        size=FS_BODY_SM,
    )

    # Desno - hijerarhijski dijagram + placeholder za sliku
    right_x = CONTENT_X + half_w + Inches(0.4)

    tree = [
        ("Kurs",       0, C_PRIMARY),
        ("Modul",      1, C_PRIMARY),
        ("Koncept",    2, C_ACCENT),
        ("Ishod",      3, C_ACCENT),
        ("Aktivnost",  4, C_MUTED),
    ]
    box_w = half_w - Inches(0.4)
    box_h = Inches(0.5)
    gap = Inches(0.12)
    y = CONTENT_Y + Inches(0.3)

    for i, (name, indent, color) in enumerate(tree):
        ind = Inches(0.4) * indent
        x = right_x + ind
        w = box_w - ind
        add_filled_rect(slide, x, y, w, box_h, fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, y, Inches(0.1), box_h, fill=color)
        add_text_box(
            slide, x + Inches(0.2), y, w - Inches(0.3), box_h,
            text=name, size=FS_BODY, bold=True, color=color,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        y += box_h + gap

    # caption ispod
    add_text_box(
        slide, right_x, y + Inches(0.1), box_w, Inches(0.5),
        text="Rekurzivna struktura, JSON serijalizacija.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )


def body_luka_prompt_template(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    # Levo - opis
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Prompt template za strukturu kursa",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Šta ulazi u prompt", "strong"),
            "Tema i kontekst kursa",
            "Lista ciljeva učenja (ishoda)",
            "Bloomov nivo koji se traži",
            "Stanje već generisanog sadržaja",
            "",
            ("Anti-redundancija", "strong"),
            "Prethodno generisane aktivnosti se ubacuju u prompt",
            "LLM se eksplicitno traži da izbegne ponavljanje",
            "Daje koherentnost preko celog kursa",
        ],
        size=FS_BODY_SM,
    )

    # Desno - skica prompta (mock code)
    right_x = CONTENT_X + half_w + Inches(0.4)
    code = (
        "Role: ti si iskusan instructional designer.\n"
        "Tema kursa: {tema}\n"
        "Ciljevi učenja:\n"
        "  - {cilj_1}\n"
        "  - {cilj_2}\n\n"
        "Generiši N aktivnosti za nivo: {bloom_level}.\n"
        "Već generisane aktivnosti (ne ponavljaj):\n"
        "  - {prev_activity_1}\n"
        "  - {prev_activity_2}\n\n"
        "Format izlaza: JSON sa poljima\n"
        "  title, description, bloom_level, estimated_time."
    )
    add_code_block(slide, right_x, CONTENT_Y, half_w, CONTENT_H - Inches(0.5),
                   code, caption="Skica prompta - parametri u {}.")


def body_uros_solo_ps4(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="SOLO Quiz Generator (Uroš)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Ulaz", "strong"),
            "PDF lekcija (Srpski / Engleski, auto-detekcija)",
            "Parsirane sekcije i learning objekti",
            "Generisana ontologija pojmova",
            "",
            ("Izlaz", "strong"),
            "MCQ pitanja na 4 SOLO nivoa",
            "Svako pitanje: source_line citat iz PDF-a",
            "Relaciona i EA: ontology anchor u tags",
            "",
            ("Tech stack", "strong"),
            "Ollama qwen2.5:14b-instruct-q4 lokalno",
            "RDFLib + SPARQL, SQLite + Flask + React",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.32),
        text="PS4 prompt template",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, right_x, CONTENT_Y + Inches(0.32), half_w, Inches(0.25),
        text="Scaria et al. (2024) - PS4 nadmašio PS5 u istom papiru",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )

    components = [
        ("Persona", "Role priming - LLM kao iskusan pedagog"),
        ("Pravila",  "SOLO def + 1 worked example (cross-domain: fotosinteza)"),
        ("Putanja",  "Chain-of-thought scaffold - Wei et al. (2022)"),
        ("Pitanja",  "JSON schema + typed distractors + Haladyna H-rules (2002)"),
    ]
    gap = Inches(0.1)
    boxes_top = CONTENT_Y + Inches(0.62)
    cite_h = Inches(0.3)
    box_h = (CONTENT_H - Inches(0.62) - 3 * gap - cite_h) / 4
    y = boxes_top

    for label, desc in components:
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.1), box_h, fill=C_ACCENT)
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.05),
            half_w - Inches(0.3), Inches(0.4),
            text=label, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.4),
            half_w - Inches(0.3), box_h - Inches(0.4),
            text=desc, size=FS_BODY_SM, color=C_TEXT,
        )
        y += box_h + gap

    # Citat ispod 4 box-a
    add_text_box(
        slide, right_x, y, half_w, cite_h,
        text="Worked example cross-domain (fotosinteza, ne OS) - Sweller & Cooper (1985):\n"
             "studira se STRUKTURA pitanja, ne sadržaj.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )


def body_two_pass_ea(slide):
    half_w = CONTENT_W / 2 - Inches(0.25)

    # Levo - Pass 1 + Pass 2
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Pass 1: pitanje + odgovor + source",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, Inches(2.4),
        [
            "LLM generiše pitanje za EA nivo",
            "Vraca tačan odgovor",
            "Vraca source_line - doslovni navod iz PDF-a",
            "Veza pitanje <-> izvor je eksplicitna",
        ],
        size=FS_BODY_SM,
    )

    y2 = CONTENT_Y + Inches(3.1)
    add_text_box(
        slide, CONTENT_X, y2, half_w, Inches(0.4),
        text="Pass 2: typed distractors (predictive prompting)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, y2 + Inches(0.5), half_w, CONTENT_H - (y2 - CONTENT_Y) - Inches(0.9),
        [
            "Echo: ponovi pitanje i tačan odgovor",
            "Generiši 3 tipizirana distraktora po SOLO nivou:",
            ("EA: APPLIES_WRONG_PRINCIPLE  /  RIGHT_PRINCIPLE_WRONG_DOMAIN  /  OVER_GENERALIZATION", "sub"),
            "Strategije zasnovane na Sadler (1998) misconception theory",
        ],
        size=FS_BODY_SM,
    )

    # Citatna napomena ispod
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + CONTENT_H - Inches(0.4), half_w, Inches(0.4),
        text="Bitew et al. (2023): ~53% distraktora ocenjeno production-ready - "
             "prevazilazi prethodne SOTA pristupe.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )

    # Desno - dijagram dva prolaza
    right_x = CONTENT_X + half_w + Inches(0.5)
    rw = CONTENT_W - half_w - Inches(0.5)

    # Pass 1 box
    box_h = Inches(2.4)
    p1 = add_filled_rect(slide, right_x, CONTENT_Y, rw, box_h,
                         fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_filled_rect(slide, right_x, CONTENT_Y, Inches(0.12), box_h, fill=C_PRIMARY)
    add_text_box(
        slide, right_x + Inches(0.25), CONTENT_Y + Inches(0.15),
        rw - Inches(0.3), Inches(0.4),
        text="PASS 1", size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
    )
    add_text_box(
        slide, right_x + Inches(0.25), CONTENT_Y + Inches(0.5),
        rw - Inches(0.3), Inches(0.5),
        text="LLM -> { pitanje, tacan_odgovor, source_line }",
        size=FS_BODY, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, right_x + Inches(0.25), CONTENT_Y + Inches(1.05),
        rw - Inches(0.3), Inches(1.2),
        text="Fokus: pedagoška valjanost pitanja i grounding u izvoru. Distraktori se NE traže u ovom prolazu.",
        size=FS_BODY_SM, color=C_TEXT,
    )

    # Strelica
    arrow_y = CONTENT_Y + box_h + Inches(0.05)
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.DOWN_ARROW, right_x + rw / 2 - Inches(0.15), arrow_y,
        Inches(0.3), Inches(0.4),
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = C_ACCENT
    arrow.line.fill.background()

    # Pass 2 box
    p2_y = arrow_y + Inches(0.5)
    p2_h = CONTENT_H - (p2_y - CONTENT_Y) - Inches(0.2)
    add_filled_rect(slide, right_x, p2_y, rw, p2_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_filled_rect(slide, right_x, p2_y, Inches(0.12), p2_h, fill=C_PRIMARY)
    add_text_box(
        slide, right_x + Inches(0.25), p2_y + Inches(0.15),
        rw - Inches(0.3), Inches(0.4),
        text="PASS 2", size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
    )
    add_text_box(
        slide, right_x + Inches(0.25), p2_y + Inches(0.5),
        rw - Inches(0.3), Inches(0.5),
        text="LLM -> { 3 typed distractor-a }",
        size=FS_BODY, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, right_x + Inches(0.25), p2_y + Inches(1.05),
        rw - Inches(0.3), p2_h - Inches(1.1),
        text="Echo prompt: LLM ponovi pitanje, pa generiše distraktore po tipiziranim strategijama greske. Smanjuje 'lažne' distraktore.",
        size=FS_BODY_SM, color=C_TEXT,
    )


def body_source_line(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Source-line citation kao anti-halucinacija",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.32), half_w, Inches(0.25),
        text="Lagana varijanta RAG-a (Lewis et al., 2020): model surface-uje svoju citaciju.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.62), half_w, CONTENT_H - Inches(0.7),
        [
            ("Šta", "strong"),
            "Uz svako pitanje, LLM mora vratiti doslovni navod iz PDF-a",
            "Navod opravdava tačan odgovor",
            "",
            ("Zašto", "strong"),
            "Ako navod ne postoji u izvoru = signal halucinacije",
            "Recenzent vidi gde je 'istina' - sekundama, ne minutima",
            "",
            ("Limit", "strong"),
            "Source_line hvata KRAĐU citata - ne hvata pogrešnu interpretaciju",
            "Za to imamo CoVe (Chain-of-Verification) - vidi sledeće slajdove",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Primer JSON izlaza",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    code = (
        '{\n'
        '  "level": "relational",\n'
        '  "question": "Kako se memorijska\\n'
        '   hijerarhija odnosi prema brzini\\n'
        '   pristupa i kapacitetu?",\n'
        '  "correct": "Brzina opada, a\\n'
        '   kapacitet raste niz hijerarhiju.",\n'
        '  "distractors": [...],\n'
        '  "source_line": "Memorijska\\n'
        '   hijerarhija je organizovana tako\\n'
        '   da brže memorije imaju manji\\n'
        '   kapacitet.",\n'
        '  "tags": { "ontology_anchor": {\\n'
        '     "source": "RAM",\\n'
        '     "target": "Cache",\\n'
        '     "type": "is_faster_than"\\n'
        '  } }\n'
        '}'
    )
    add_code_block(slide, right_x, CONTENT_Y + Inches(0.5), half_w,
                   CONTENT_H - Inches(0.6), code,
                   caption="Pojednostavljeni primer odgovora LLM-a.")


def body_stefan_dekompozicija(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Computer Use Video Instructions (Stefan)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Ulaz", "strong"),
            "Prirodno-jezička instrukcija praktičkog zadatka",
            "Npr. \"Kreiraj C# konzolni projekat u Visual Studio-u\"",
            "",
            ("Izlaz", "strong"),
            "Snimljeni MP4 video tutorijal",
            "OWL fajl izvršnog plana (auditable)",
            "",
            ("Hijerarhija razgradnje", "strong"),
            "Task: kompletan zadatak (zatvorena celina)",
            "Step: atomična akcija (otvori, klikni, otkucaj)",
            "Action: tip akcije iz fiksne ontologije akcija",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Task / Step / Action hijerarhija",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    # 3 nivoa kartice
    levels = [
        ("Task", "Open VS, create C# Console app, write Hello World"),
        ("Step", "1. open_application Visual Studio   ->   2. wait   ->   3. click Create new project   ->   ..."),
        ("Action", "open_application | click | type_text | key_press | wait | capture_screen"),
    ]
    gap = Inches(0.15)
    box_h = (CONTENT_H - Inches(0.5) - 2 * gap) / 3
    y = CONTENT_Y + Inches(0.5)

    for i, (lbl, content) in enumerate(levels):
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.1), box_h, fill=C_ACCENT)
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.1),
            half_w - Inches(0.3), Inches(0.35),
            text=lbl, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.45),
            half_w - Inches(0.3), box_h - Inches(0.45),
            text=content, size=FS_BODY_SM, color=C_TEXT,
        )
        y += box_h + gap


def body_prompt_vs_finetune(slide):
    headers = ["", "Prompt engineering", "Fine-tuning"]
    rows = [
        ["Trosak razvoja",     "Nizak - samo prompt",                 "Visok - dataset + trening"],
        ["Trosak izvršenja",   "Veci prompt = veći token cost",       "Manji prompt, ali specijalizovan model"],
        ["Brzina iteracije",   "Sekunde - promenis prompt",           "Sati / dani - rerun trening"],
        ["Domen-specifičnost", "Ogranicena, oslanja se na bazni model","Visoka, prilagođen domenu"],
        ["Kontrola halucinacija","Posredna (CoT, schema, grounding)","Direktna (kvalitet treninga)"],
        ["Nas izbor",          "Da - PS4 + ontology grounding",       "Ne, za sada"],
    ]
    add_styled_table(slide, CONTENT_X, CONTENT_Y, CONTENT_W,
                     CONTENT_H - Inches(1.4), headers, rows,
                     col_widths=[1.2, 1.7, 1.7])

    # Donja poruka
    y2 = CONTENT_Y + CONTENT_H - Inches(1.2)
    add_filled_rect(slide, CONTENT_X, y2, CONTENT_W, Inches(1.0),
                    fill=C_CARD_BG, line=None)
    add_filled_rect(slide, CONTENT_X, y2, Inches(0.08), Inches(1.0), fill=C_ACCENT)
    add_text_box(
        slide, CONTENT_X + Inches(0.25), y2 + Inches(0.1),
        CONTENT_W - Inches(0.4), Inches(0.4),
        text="Zašto prompt engineering, ne fine-tune?", size=FS_BODY, bold=True,
        color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X + Inches(0.25), y2 + Inches(0.45),
        CONTENT_W - Inches(0.4), Inches(0.6),
        text="Bez stabilnog domenskog dataset-a za SOLO/Bloom u programiranju, prompt + ontology grounding daje brzu, transparentnu i jeftiniju iteraciju.",
        size=FS_BODY_SM, color=C_TEXT,
    )


def body_local_vs_cloud(slide):
    headers = ["", "Lokalni: Ollama Qwen 2.5 14B", "Cloud: Groq Llama 3.3 70B", "Cloud Vision: Qwen 2.5 VL 72B"]
    rows = [
        ["Koristi se za", "Uroš - SOLO pitanja", "Stefan - planiranje koraka", "Stefan - UI detekcija"],
        ["Trosak po pozivu", "0 (lokalno)", "Mali (per token)", "Srednji (vision)"],
        ["Brzina", "Sporije, zavisi od GPU-a", "Vrlo brzo (Groq accel.)", "Brza"],
        ["Privatnost",  "Visoka (sve lokalno)", "Niska - cloud", "Niska - cloud"],
        ["Kvalitet",    "Solidan za EA, slabiji za nuansiranja", "Veoma dobar", "Najbolji za vision"],
        ["Cache",       "SHA-256 keys u SQLite",  "n/a u ovom projektu", "n/a"],
    ]
    add_styled_table(slide, CONTENT_X, CONTENT_Y, CONTENT_W,
                     CONTENT_H - Inches(0.4), headers, rows,
                     col_widths=[1.0, 1.6, 1.6, 1.6])


def body_seg_a_kljucna(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), CONTENT_W, Inches(0.4),
        text="Ključna poruka segmenta A",
        size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )

    box_w = CONTENT_W * 0.85
    box_x = CONTENT_X + (CONTENT_W - box_w) / 2
    box_y = CONTENT_Y + Inches(1.1)
    box_h = Inches(2.6)

    add_filled_rect(slide, box_x, box_y, box_w, box_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_horizontal_line(slide, box_x, box_y, box_w, color=C_PRIMARY,
                        height=Emu(30000))

    add_text_box(
        slide, box_x + Inches(0.4), box_y + Inches(0.5), box_w - Inches(0.8),
        Inches(1.6),
        text=("Hijerarhijska dekompozicija + formalna pedagoška "
              "taksonomija ugrađena u prompt transformisu LLM iz "
              "generickog generatora teksta u pedagoški kalibrisan alat."),
        size=22, italic=True, color=C_PRIMARY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    # ispod - 3 mini-poente
    yy = box_y + box_h + Inches(0.4)
    pts = [
        ("Bloom",       "drives the WHAT - tipovi aktivnosti"),
        ("SOLO",        "drives the HOW DEEP - kognitivne razlike"),
        ("PS4 / dekompozicija", "drives the HOW - struktura prompta"),
    ]
    gap = Inches(0.2)
    col_w = (CONTENT_W - 2 * gap) / 3
    for i, (h, b) in enumerate(pts):
        x = CONTENT_X + i * (col_w + gap)
        add_text_box(slide, x, yy, col_w, Inches(0.4),
                     text=h, size=FS_BODY, bold=True, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, x, yy + Inches(0.4), col_w, Inches(0.6),
                     text=b, size=FS_BODY_SM, color=C_MUTED,
                     align=PP_ALIGN.CENTER)


# --- SEGMENT B helperi -------------------------------------------------------

def body_owl_uvod(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="OWL (Web Ontology Language)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "W3C standard za formalne ontologije",
            "Definise klase, instance, relacije, ograničenja",
            "Čitljiv mas i mas (RDF/XML, Turtle, OWL/XML)",
            "Reasoning - može da izvodi nove činjenice",
            "",
            ("Mi koristimo", "strong"),
            "Stefan: OWL u Turtle (computer_use.ttl)",
            "Uroš: seed OWL + RDF/XML eksport za Protege",
            "Luka: konceptualni model serijalizovan u JSON",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    code = (
        "@prefix cu: <http://example.org/computer-use#> .\n"
        "@prefix owl: <.../owl#> .\n"
        "@prefix rdf: <.../rdf-syntax-ns#> .\n\n"
        "cu:Task a owl:Class .\n"
        "cu:Step a owl:Class .\n"
        "cu:hasStep a owl:ObjectProperty ;\n"
        "    rdfs:domain cu:Task ;\n"
        "    rdfs:range  cu:Step .\n"
        "cu:stepOrder a owl:DatatypeProperty ;\n"
        "    rdfs:domain cu:Step ;\n"
        "    rdfs:range  xsd:integer ."
    )
    add_code_block(slide, right_x, CONTENT_Y + Inches(0.5), half_w,
                   CONTENT_H - Inches(0.6), code,
                   caption="Mini-isečak iz computer_use.ttl (Stefan).")


def body_rdf_sparql(slide):
    # Levo - RDF
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="RDF - tripleti",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, Inches(2.5),
        [
            "Resource Description Framework",
            "Sve činjenice = (subject, predicate, object)",
            "Primer: (RAM, is_faster_than, Disk)",
            "OWL je TBox, RDF je ABox - klase vs instance",
        ],
        size=FS_BODY_SM,
    )

    code_rdf = (
        ":RAM   rdf:type  :MemoryDevice .\n"
        ":Disk  rdf:type  :MemoryDevice .\n"
        ":RAM   :isFasterThan  :Disk ."
    )
    add_code_block(
        slide, CONTENT_X, CONTENT_Y + Inches(3.1),
        half_w, Inches(1.3), code_rdf,
        caption="3 RDF tripleta = mini graf.",
    )

    # Desno - SPARQL
    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="SPARQL - upitni jezik nad RDF/OWL",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    code_sparql = (
        "SELECT ?step ?order ?action ?target WHERE {\n"
        "  ?task a cu:Task ; cu:hasStep ?step .\n"
        "  ?step cu:stepOrder ?order ;\n"
        "        cu:hasAction ?action ;\n"
        "        cu:targetName ?target .\n"
        "} ORDER BY ?order"
    )
    add_code_block(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w,
        CONTENT_H - Inches(0.7), code_sparql,
        caption="Stefanov upit za čitanje koraka iz OWL plana.",
    )


def body_luka_json_primer(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="JSON serijalizacija strukture kursa",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "Hijerarhijska, rekurzivna struktura",
            "Svaki čvor ima id, name, bloom_level, children",
            "Aktivnosti su listovi u stablu",
            "Lako se mapira u Protege / OWL graph",
            "",
            ("Zašto JSON, a ne odmah OWL?", "strong"),
            "Lakše za LLM da generiše validan JSON",
            "Lakše za UI prikaz i editovanje",
            "Mapiranje JSON -> OWL je determinizam, lako se radi unazad",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    code_json = (
        '{\n'
        '  "course": "Operativni sistemi",\n'
        '  "modules": [\n'
        '    {\n'
        '      "name": "Procesi i niti",\n'
        '      "concepts": [\n'
        '        {\n'
        '          "name": "Proces",\n'
        '          "outcomes": [\n'
        '            {\n'
        '              "name": "Razumeti stanja procesa",\n'
        '              "bloom": "Razumeti",\n'
        '              "activities": [...]\n'
        '            }\n'
        '          ]\n'
        '        }\n'
        '      ]\n'
        '    }\n'
        '  ]\n'
        '}'
    )
    add_code_block(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w,
        CONTENT_H - Inches(0.6), code_json,
        caption="Primer Lukinog JSON izlaza (pojednostavljeno).",
    )


def body_uros_ontologija(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="ConceptRelationship ontologija",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Modeluje", "strong"),
            "Relacije između learning objekata (LO) jedne lekcije",
            "Edge = (source, target, type, description)",
            "Čuva se u SQL tabeli i eksportuje u Turtle / OWL",
            "",
            ("Multi-pass ekstrakcija", "strong"),
            "5 LLM prolaza, jedan po tipu relacije",
            "Smanjuje hallucinated 'prerequisite' lance",
            "Konzervativan fallback samo ako LLM vrati nulu",
            "",
            ("Tipovi relacija", "strong"),
            "prerequisite, example_of, depends_on, ...",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Tipovi relacija u ontologiji",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    rels = [
        ("prerequisite",  "A se mora razumeti pre B"),
        ("example_of",    "A je primer pojma B"),
        ("depends_on",    "A operativno zavisi od B"),
        ("similar_to",    "A i B dele ključnu karakteristiku"),
        ("part_of",       "A je deo strukture B"),
    ]
    gap = Inches(0.1)
    n = len(rels)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)
    for name, desc in rels:
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.1), box_h, fill=C_ACCENT)
        add_text_box(
            slide, right_x + Inches(0.25), y, Inches(2.0), box_h,
            text=name, size=FS_BODY, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, right_x + Inches(2.3), y, half_w - Inches(2.4), box_h,
            text=desc, size=FS_BODY_SM, color=C_MUTED,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        y += box_h + gap


def body_ontology_anchor(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Ontology anchor kao kontekst pitanja",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.32), half_w, Inches(0.25),
        text="KAQG pattern (Lin et al., 2025) - knowledge graph kao gradnja-constraint.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.62), half_w, CONTENT_H - Inches(0.7),
        [
            "Za relacionalna i EA pitanja",
            "SPARQL bira jedan ConceptRelationship red",
            "Anchor postaje deo prompta",
            "Pitanje se gradi oko jedne konkretne veze",
            "",
            ("Posledice", "strong"),
            "Pitanje je trasabilno - zna se na koju vezu se odnosi",
            "Dedup po (anchor, normalized_correct_answer) - egzaktno",
            "Lakše za reviziju nastavnika",
            "",
            ("Bez anchor-a", "strong"),
            "Relacionalna pitanja 'klize' u genericke poređenja",
            "Anchor pretvara vagi nivo u specifičnu ivicu grafa",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    code = (
        "SELECT ?src ?type ?tgt ?desc WHERE {\n"
        "  ?rel a :ConceptRelationship ;\n"
        "       :source ?src ; :target ?tgt ;\n"
        "       :relType ?type ; :description ?desc .\n"
        "  ?src :inLesson <lesson_42> .\n"
        "} ORDER BY RAND() LIMIT 1"
    )
    add_code_block(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w,
        CONTENT_H - Inches(0.6), code,
        caption="Biranje jedne veze kao anchor-a za prompt LLM-a.",
    )


def body_stefan_owl(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Stefan - OWL ontologija plana",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Klase", "strong"),
            "Task - kompletan automatizacioni zadatak",
            "Step - atomična akcija u zadatku",
            "Action - tip akcije iz fiksne liste",
            "UIElement, Application, ExecutionState",
            "",
            ("Object properties", "strong"),
            "hasStep, hasAction, hasTarget",
            "nextStep / previousStep, requiresApplication",
            "",
            ("Data properties", "strong"),
            "stepOrder, stepDescription, inputValue",
            "targetName, waitDuration, expectedResult",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Primer instance Step-a u OWL/RDF",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    code = (
        "<cu:Step rdf:about=\".../Step_1\">\n"
        "  <cu:stepOrder>1</cu:stepOrder>\n"
        "  <cu:hasAction rdf:resource=\"...open_application\"/>\n"
        "  <cu:targetName>Visual Studio</cu:targetName>\n"
        "  <cu:stepDescription>\n"
        "    Start Visual Studio\n"
        "  </cu:stepDescription>\n"
        "  <cu:expectedResult>\n"
        "    Visual Studio is opened\n"
        "  </cu:expectedResult>\n"
        "  <cu:hasState\n"
        "     rdf:resource=\"...PendingState\"/>\n"
        "</cu:Step>"
    )
    add_code_block(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w,
        CONTENT_H - Inches(0.6), code,
        caption="Skiniran iz computer_use.ttl (Stefanov repo).",
    )


def body_validacija_pravila(slide):
    headers = ["Pravilo", "Opis", "Zašto"]
    rows = [
        ["Action validation",
         "Samo akcije iz ontologije su dozvoljene",
         "LLM ne može izmisliti akciju koja ne može biti izvrsena"],
        ["Sequence validation",
         "open_application mora prethoditi interakciji",
         "Sprecava klikove na nepostojeći prozor"],
        ["Wait recommendation",
         "wait posle open_application",
         "UI se ucitava - bez pauze klikovi promasuju"],
        ["Parameter validation",
         "type_text zahteva value, key_press zahteva ime tipke",
         "Hvata strukturne greske LLM-a u JSON izlazu"],
        ["Target validation",
         "click mora imati specifičan target",
         "Bez targeta nema šta da se klikne"],
    ]
    add_styled_table(slide, CONTENT_X, CONTENT_Y, CONTENT_W,
                     CONTENT_H - Inches(0.4), headers, rows,
                     col_widths=[1.2, 2.2, 2.2])


def body_anatomija_halucinacije(slide):
    """Dva primera jedan-do-drugog: lose (bez source_line) vs dobro (sa check-om)."""
    half_w = CONTENT_W / 2 - Inches(0.2)

    # Leva strana - LOSE: pitanje bez source_line check-a
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Bez source_line check-a", size=FS_SUBTITLE, bold=True,
        color=C_ACCENT,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), half_w, Inches(0.3),
        text="(LLM može izmisliti činjenicu)",
        size=FS_BODY_SM, italic=True, color=C_MUTED,
    )

    bad_y = CONTENT_Y + Inches(0.8)
    bad_h = CONTENT_H - Inches(1.0)
    add_filled_rect(slide, CONTENT_X, bad_y, half_w, bad_h,
                    fill=C_CARD_BG, line=C_ACCENT, line_w=1.2)
    add_filled_rect(slide, CONTENT_X, bad_y, half_w, Inches(0.08), fill=C_ACCENT)

    pad = Inches(0.22)
    inner_x = CONTENT_X + pad
    inner_w = half_w - 2 * pad
    cur_y = bad_y + Inches(0.25)

    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.3),
                 text="Q: Koja godina je donesen TCP/IP standard?",
                 size=FS_BODY_SM, bold=True, color=C_PRIMARY)
    cur_y += Inches(0.55)
    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.3),
                 text="A: 1981. (RFC 791)",
                 size=FS_BODY_SM, italic=True, color=C_TEXT)
    cur_y += Inches(0.55)
    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.3),
                 text="source_line: (nedostaje)",
                 size=FS_BODY_SM, color=C_MUTED)
    cur_y += Inches(0.5)
    add_filled_rect(slide, inner_x, cur_y, inner_w, Inches(1.0),
                    fill=C_WHITE, line=C_ACCENT, line_w=0.75)
    add_text_box(slide, inner_x + Inches(0.15), cur_y + Inches(0.12),
                 inner_w - Inches(0.3), Inches(0.3),
                 text="! Problem", size=FS_BODY_SM, bold=True, color=C_ACCENT)
    add_text_box(
        slide, inner_x + Inches(0.15), cur_y + Inches(0.4),
        inner_w - Inches(0.3), Inches(0.6),
        text="Tačan podatak izgleda - ali je u PDF-u pisalo '1983 (DoD)'. "
             "LLM se 'setio' opšte činjenice umesto da čita lekciju.",
        size=FS_BODY_SM, color=C_TEXT,
    )

    # Desna strana - DOBRO: sa source_line + fuzzy match check
    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Sa source_line + fuzzy match", size=FS_SUBTITLE, bold=True,
        color=C_PRIMARY,
    )
    add_text_box(
        slide, right_x, CONTENT_Y + Inches(0.4), half_w, Inches(0.3),
        text="(halucinacija se markira pre nego stigne do nastavnika)",
        size=FS_BODY_SM, italic=True, color=C_MUTED,
    )

    add_filled_rect(slide, right_x, bad_y, half_w, bad_h,
                    fill=C_CARD_BG, line=C_PRIMARY, line_w=1.2)
    add_filled_rect(slide, right_x, bad_y, half_w, Inches(0.08), fill=C_PRIMARY)

    inner_x = right_x + pad
    inner_w = half_w - 2 * pad
    cur_y = bad_y + Inches(0.25)

    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.3),
                 text="Q: Koja godina je donesen TCP/IP standard?",
                 size=FS_BODY_SM, bold=True, color=C_PRIMARY)
    cur_y += Inches(0.55)
    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.3),
                 text="A: 1983.",
                 size=FS_BODY_SM, italic=True, color=C_TEXT)
    cur_y += Inches(0.55)
    add_text_box(slide, inner_x, cur_y, inner_w, Inches(0.6),
                 text="source_line: \"1. januara 1983. ARPANET je presao na "
                      "TCP/IP protokol stek.\"",
                 size=FS_BODY_SM, color=C_TEXT)
    cur_y += Inches(0.75)
    add_filled_rect(slide, inner_x, cur_y, inner_w, Inches(1.0),
                    fill=C_WHITE, line=C_PRIMARY, line_w=0.75)
    add_text_box(slide, inner_x + Inches(0.15), cur_y + Inches(0.12),
                 inner_w - Inches(0.3), Inches(0.3),
                 text="OK Fuzzy match: 0.94", size=FS_BODY_SM, bold=True,
                 color=C_PRIMARY)
    add_text_box(
        slide, inner_x + Inches(0.15), cur_y + Inches(0.4),
        inner_w - Inches(0.3), Inches(0.6),
        text="Navod nađen u PDF-u, sekcija 2.1, str. 14. Pitanje prolazi "
             "automatsku validaciju i ide u QuestionBank.",
        size=FS_BODY_SM, color=C_TEXT,
    )

    # Footer message
    y_msg = CONTENT_Y + CONTENT_H - Inches(0.18)
    add_text_box(
        slide, CONTENT_X, y_msg, CONTENT_W, Inches(0.3),
        text="Pal et al. (2023): bez RAG/grounding-a, GPT-4 ima 15-20% stopu halucinacije na faktičkim pitanjima.",
        size=FS_CAPTION, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER,
    )


def body_anti_halucinacije(slide):
    cols = [
        ("Uroš",
         "source_line citat",
         ["LLM mora vratiti doslovni navod iz PDF-a",
          "Citat opravdava tačan odgovor",
          "Ako navod ne postoji = signal halucinacije",
          "Recenzent vidi gde je tačan odgovor 'usidren'"]),
        ("Luka",
         "state-of-generation kontekst",
         ["Već generisane aktivnosti ubacuju se u prompt",
          "LLM eksplicitno čita 'ne ponavljaj ovo'",
          "Daje koherentnost kroz ceo kurs",
          "Smanjuje generičke, ponavljajuce ishode"]),
        ("Stefan",
         "SPARQL odbacuje nevalidne sekvence",
         ["Pre izvršavanja - ontološka provera",
          "Pravila: open_app pre interakcije, wait posle open_app",
          "type_text zahteva value, ...",
          "Nevalidan plan = stop, ne pravi se video"]),
    ]
    gap = Inches(0.25)
    col_w = (CONTENT_W - 2 * gap) / 3

    for i, (who, title, points) in enumerate(cols):
        x = CONTENT_X + i * (col_w + gap)
        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.1),
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.2)
        cur_y = CONTENT_Y + Inches(0.25)
        add_text_box(slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.3),
                     text=who.upper(), size=FS_SECTION_LABEL, bold=True,
                     color=C_ACCENT)
        cur_y += Inches(0.32)
        add_text_box(slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.5),
                     text=title, size=20, bold=True, color=C_PRIMARY)
        cur_y += Inches(0.55)
        add_horizontal_line(slide, x + pad, cur_y, Inches(0.5),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.15)
        add_bullet_list(slide, x + pad, cur_y, col_w - 2 * pad,
                        CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.3),
                        points, size=FS_BODY_SM)


def body_seg_b_kljucna(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), CONTENT_W, Inches(0.4),
        text="Ključna poruka segmenta B",
        size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    box_w = CONTENT_W * 0.85
    box_x = CONTENT_X + (CONTENT_W - box_w) / 2
    box_y = CONTENT_Y + Inches(1.1)
    box_h = Inches(2.6)
    add_filled_rect(slide, box_x, box_y, box_w, box_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_horizontal_line(slide, box_x, box_y, box_w, color=C_PRIMARY,
                        height=Emu(30000))
    add_text_box(
        slide, box_x + Inches(0.4), box_y + Inches(0.4), box_w - Inches(0.8),
        box_h - Inches(0.8),
        text=("Ontologija je ugovor između LLM-a i sistema: ono sto LLM "
              "proizvede mora se mapirati na klase i relacije. "
              "To filtrira halucinacije i omogućava ponovnu upotrebu sadržaja."),
        size=22, italic=True, color=C_PRIMARY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )

    yy = box_y + box_h + Inches(0.4)
    pts = [
        ("OWL/RDF",  "definicija strukture i instanci"),
        ("SPARQL",   "kontrola toka i upita"),
        ("Anchor",   "vezivanje pitanja za izvor"),
    ]
    gap = Inches(0.2)
    col_w = (CONTENT_W - 2 * gap) / 3
    for i, (h, b) in enumerate(pts):
        x = CONTENT_X + i * (col_w + gap)
        add_text_box(slide, x, yy, col_w, Inches(0.4),
                     text=h, size=FS_BODY, bold=True, color=C_ACCENT,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, x, yy + Inches(0.4), col_w, Inches(0.6),
                     text=b, size=FS_BODY_SM, color=C_MUTED,
                     align=PP_ALIGN.CENTER)


# --- SEGMENT C helperi -------------------------------------------------------

def body_validation_stack(slide):
    """3 sloja a-priori validacije: prevention + generation + detection."""
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="Pre nego sto bilo koji student vidi pitanje - 3 sloja provere",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), CONTENT_W, Inches(0.3),
        text="Isti rule kodovi (H14, H17, H22, H27, ...) postoje i u promptu (prevent) i u lint-u (detect).",
        size=FS_BODY_SM, italic=True, color=C_MUTED,
    )

    layers = [
        ("1. PREVENTION", "Pre LLM poziva", C_PRIMARY,
         ["PS4 template (Scaria 2024) - role + def + 1 worked example + CoT + schema",
          "Haladyna 2002 rules in prompt: S1-S4 (stem), O1-O7 (options)",
          "Typed distractor strategies po SOLO nivou (Bitew 2023 + Sadler 1998)",
          "Ontology anchor block za relacionalna/EA pitanja (KAQG 2025)"]),
        ("2. GENERATION", "LLM poziv + source grounding", C_ACCENT,
         ["Ollama Qwen 2.5 14B, single call, JSON mode",
          "Verbatim source_line citacija (lagani RAG - Lewis 2020)",
          "Two-pass orchestration za Extended Abstract (Bitew 2023)",
          "Self-consistency best-of-N (Wang 2022) za high-stakes generisanje"]),
        ("3. DETECTION", "Posle LLM-a, pre čoveka", C_PRIMARY,
         ["Haladyna lint (11 pravila, 0-100 skor) + Embedding plausibility/diversity",
          "Chain-of-Verification (Dhuliawala 2023) - LLM proverava sebe",
          "Solvability test (LLM kao slep solver) + SOLO judge (Cohen's kappa)",
          "Concept coverage v2 (Kurdi 2020) + IOC + readability + ambiguity"]),
    ]

    n = len(layers)
    gap = Inches(0.15)
    row_h = (CONTENT_H - Inches(0.85) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.85)

    for label, sub, color, points in layers:
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, row_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.15), row_h, fill=color)

        # Levo - label
        add_text_box(
            slide, CONTENT_X + Inches(0.3), y + Inches(0.08),
            Inches(2.6), Inches(0.35),
            text=label, size=FS_BODY, bold=True, color=color,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.3), y + Inches(0.42),
            Inches(2.6), Inches(0.3),
            text=sub, size=FS_BODY_SM, italic=True, color=C_MUTED,
        )

        # Desno - bullet lista tehnika
        add_bullet_list(
            slide, CONTENT_X + Inches(3.05), y + Inches(0.08),
            CONTENT_W - Inches(3.2), row_h - Inches(0.15),
            points, size=FS_BODY_SM,
        )

        y += row_h + gap


def body_haladyna_lint(slide):
    """11 automatizovanih Haladyna pravila - prevent + detect hybrid."""
    half_w = CONTENT_W / 2 - Inches(0.2)

    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Haladyna et al. (2002) - 31 pravilo, 11 automatizovanih",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.32), half_w, Inches(0.25),
        text="Validovano nad 27 udžbenika i 27 empirijskih studija.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )

    # Levo - tabela glavnih H-pravila
    headers = ["Kod", "Pravilo"]
    rows = [
        ["H14", "Stem se završava pitanjem ili jasnim imperativom"],
        ["H16", "Stem ispod ~250 znakova, bez ukrasa"],
        ["H17", "Izbeći negaciju (ili je istakni)"],
        ["H19", "Tačno jedna tačna opcija"],
        ["H22", "Nijedne dve opcije nisu parafraze"],
        ["H24", "Najduža opcija <= 2x najkraća (length clue)"],
        ["H25", "Nikad 'all/none of the above'"],
        ["H27", "Tačna opcija NIJE najduža (give-away)"],
        ["O7",  "Sve opcije gramatički paralelne"],
    ]
    add_styled_table(
        slide, CONTENT_X, CONTENT_Y + Inches(0.65), half_w,
        CONTENT_H - Inches(0.75), headers, rows,
        col_widths=[1, 4.5],
    )

    # Desno - hybrid prevent+detect dijagram
    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Hybrid: ista pravila u promptu I u lint-u",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    # 3 box-a sa strelicama nadole
    flow = [
        ("PROMPT (prevent)", "STEM_RULES + OPTION_RULES sekcije",
         "LLM eksplicitno čita 11 numerisanih ograničenja", C_PRIMARY),
        ("LLM CALL", "Ollama Qwen 2.5 14B",
         "Generiše JSON sa stem, correct, distractors, source_line", C_ACCENT),
        ("LINT (detect)", "mcq_lint.py rule engine",
         "Skenira izlaz, vraća H-flag-ove + 0-100 composite skor", C_PRIMARY),
    ]
    n = len(flow)
    arrow_h = Inches(0.3)
    boxes_top = CONTENT_Y + Inches(0.55)
    avail = CONTENT_H - Inches(1.4)
    box_h = (avail - (n - 1) * arrow_h) / n
    y = boxes_top

    for i, (lbl, sub, desc, color) in enumerate(flow):
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.1), box_h, fill=color)
        add_text_box(
            slide, right_x + Inches(0.22), y + Inches(0.08),
            half_w - Inches(0.3), Inches(0.3),
            text=lbl, size=FS_BODY_SM, bold=True, color=color,
        )
        add_text_box(
            slide, right_x + Inches(0.22), y + Inches(0.38),
            half_w - Inches(0.3), Inches(0.25),
            text=sub, size=FS_CAPTION, italic=True, color=C_MUTED,
        )
        add_text_box(
            slide, right_x + Inches(0.22), y + Inches(0.62),
            half_w - Inches(0.3), box_h - Inches(0.65),
            text=desc, size=FS_BODY_SM, color=C_TEXT,
        )

        if i < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                right_x + half_w / 2 - Inches(0.15), y + box_h + Inches(0.02),
                Inches(0.3), arrow_h - Inches(0.04),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C_ACCENT
            arrow.line.fill.background()

        y += box_h + arrow_h

    # Donja poruka
    y_msg = CONTENT_Y + CONTENT_H - Inches(0.4)
    add_text_box(
        slide, right_x, y_msg, half_w, Inches(0.4),
        text="Isti rule kod (H22) i u promptu i u lint-u - "
             "lint flag se može kvotirati nazad u prompt rule.",
        size=FS_CAPTION, italic=True, color=C_ACCENT, align=PP_ALIGN.CENTER,
    )


def body_llm_kao_kriticar(slide):
    """LLM kao samokritičar: CoVe + Solvability + SOLO judge + IOC."""
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="LLM kao kritičar samog sebe - 4 tehnike",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), CONTENT_W, Inches(0.3),
        text="Drugi LLM (ili isti sa različitim promptom + temp=0.1) procenjuje izlaz prvog.",
        size=FS_BODY_SM, italic=True, color=C_MUTED,
    )

    cards = [
        ("CoVe",
         "Chain-of-Verification",
         "Dhuliawala et al. (2023)",
         [
             "1. uzmi pitanje + tačan odgovor",
             "2. planiraj 2-3 kratka verifikacijska pitanja",
             "3. odgovori NA SVAKO nezavisno (samo iz izvora)",
             "4. sudi: SUPPORTED / UNDERDETERMINED / CONTRADICTED",
             "Hvata slučaj kad citat postoji, ali NE opravdava odgovor",
         ]),
        ("Solvability",
         "LLM kao slep solver",
         "Crocker & Algina (1986) p-value",
         [
             "Sakrij key, izmesaj redosled opcija, N=5 prolaza",
             "p ~ 1.0 -> trivijalno (ili length clue)",
             "0.6-0.9 -> primerena težina za LLM-class solver",
             "p < 0.5 -> misframed / ključ pogrešan / out of source",
             "Sintetički analog klasičnog test-theory difficulty indeksa",
         ]),
        ("SOLO judge",
         "Nezavisna klasifikacija + Cohen's kappa",
         "Landis & Koch (1977)",
         [
             "Drugi LLM klasifikuje pitanje u SOLO nivo",
             "Bez znanja koji je nivo generator trazio",
             "kappa < 0.20 slight | 0.41-0.60 moderate",
             "0.61-0.80 substantial | 0.81+ almost perfect",
             "A-priori provera 'da li pogođeno SOLO' bez studenata",
         ]),
        ("IOC",
         "Item-Objective Congruence",
         "Rovinelli & Hambleton (1977)",
         [
             "Za svaki item: ocena -1 / 0 / +1 vs ishod učenja",
             "Da li ovo pitanje meri BAS ovaj ishod?",
             "IOC index = srednja vrednost, [-1, +1]",
             ">= 0.5 prihvatljivo, >= 0.75 jako",
             "Content validity - prvi sloj psihometrije",
         ]),
    ]

    gap = Inches(0.18)
    col_w = (CONTENT_W - gap) / 2
    row_top = CONTENT_Y + Inches(0.78)
    row_h = (CONTENT_H - Inches(0.85) - gap) / 2

    for i, (name, sub, ref, points) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = CONTENT_X + col * (col_w + gap)
        y = row_top + row * (row_h + gap)

        add_filled_rect(slide, x, y, col_w, row_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.22)
        cur_y = y + Inches(0.18)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.32),
            text=name, size=20, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, x + col_w - Inches(2.5) - pad, cur_y,
            Inches(2.5), Inches(0.32),
            text=ref, size=FS_CAPTION, italic=True, color=C_MUTED,
            align=PP_ALIGN.RIGHT,
        )
        cur_y += Inches(0.35)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.28),
            text=sub, size=FS_BODY_SM, italic=True, color=C_ACCENT,
        )
        cur_y += Inches(0.3)
        add_horizontal_line(slide, x + pad, cur_y, Inches(0.4),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.12)
        add_bullet_list(
            slide, x + pad, cur_y, col_w - 2 * pad,
            row_h - (cur_y - y) - Inches(0.15),
            points, size=FS_BODY_SM, line_spacing=1.15, space_after=2,
        )


def body_pilot_metodologija(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Pilot evaluacija SOLO pitanja (Uroš)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Cilj evaluacije", "strong"),
            "Da li je kvalitet konstantan kroz SOLO nivoe?",
            "Da li distraktori postaju lažni / slabi na višim nivoima?",
            "",
            ("Metodologija", "strong"),
            "Više PDF lekcija iz različitih domena",
            "N pitanja po nivou, ručna anotacija",
            "Sudije: članovi tima + nezavisni recenzent",
            "Metrike: valjanost, kvalitet distraktora, source-line tačnost",
            "",
            ("Format izlaza", "strong"),
            "Skor na skali 1-5 po dimenziji",
            "Komentari za kvalitativnu analizu",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_image_placeholder(
        slide, right_x, CONTENT_Y, half_w, CONTENT_H,
        label="Tabela metrike po nivou / domenu",
    )


def body_kvalitet_po_nivou(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Rezultati - kvalitet po SOLO nivou",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    # Mali "graf" simuliran preko bar-charta od pravougaonika
    levels = [("Unistructural", 0.92),
              ("Multistructural", 0.85),
              ("Relational", 0.70),
              ("Extended Abstract", 0.55)]
    gap = Inches(0.18)
    n = len(levels)
    chart_y = CONTENT_Y + Inches(0.5)
    chart_h = CONTENT_H - Inches(1.2)
    row_h = (chart_h - (n - 1) * gap) / n

    max_bar_w = half_w - Inches(1.6)

    for i, (lvl, val) in enumerate(levels):
        y = chart_y + i * (row_h + gap)
        add_text_box(
            slide, CONTENT_X, y, Inches(1.5), row_h,
            text=lvl, size=FS_BODY_SM, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        # bar
        bar_w = max_bar_w * val
        bx = CONTENT_X + Inches(1.55)
        add_filled_rect(slide, bx, y + row_h * 0.15, max_bar_w,
                        row_h * 0.7, fill=C_LIGHT_LINE, line=None)
        add_filled_rect(slide, bx, y + row_h * 0.15, bar_w,
                        row_h * 0.7, fill=C_ACCENT, line=None)
        # value
        add_text_box(
            slide, bx + bar_w + Inches(0.1), y, Inches(0.7), row_h,
            text=f"{int(val*100)}%", size=FS_BODY, bold=True,
            color=C_PRIMARY, anchor=MSO_ANCHOR.MIDDLE,
        )

    # Caption ispod
    add_text_box(
        slide, CONTENT_X, chart_y + chart_h + Inches(0.1), half_w, Inches(0.5),
        text="Orijentaciono, N=80 anotiranih pitanja iz 4 lekcije OS-a "
             "(2 člana tima + 1 nezavisni recenzent). Brojke se prosiruju u publikaciji.",
        size=FS_CAPTION, italic=True, color=C_MUTED,
    )

    # Desno - šta to znači
    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Šta vidimo u rezultatima",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "Validnost odgovora pada sa porastom SOLO nivoa",
            "Distraktori postaju manje uverljivi",
            "EA: često se 'distraktor' može odbraniti kao tačan",
            "",
            ("Strukturna posledica - nije bug", "strong"),
            "Lister et al. (2006): EA zadaci su strukturalno teško",
            "i za STUDENTE programiranja - nije iznenađenje da i LLM",
            "tu zaostaje. Razlog je isti: transfer između konteksta.",
            "",
            ("Posledica za sistem", "strong"),
            "Two-pass + typed distractors pomazu, ne rešavaju potpuno",
            "Ontology anchor stabilizuje relacionalna pitanja",
            "Na EA nivou - obavezna recenzija nastavnika",
        ],
        size=FS_BODY_SM,
    )


def body_primer_solo_pitanja(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="Primeri pitanja po SOLO nivoima",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    levels = [
        ("Unistructural",
         "Koja struktura procesa čuva PID?",
         "PCB - Process Control Block"),
        ("Multistructural",
         "Koja tri stanja procesa najčešće prepoznajemo?",
         "Ready, Running, Blocked"),
        ("Relational",
         "Kako kontekstno prebacivanje utiče na throughput sistema?",
         "Povecava overhead, smanjuje koristan rad CPU-a."),
        ("Extended Abstract",
         "Kako biste prilagodili scheduling algoritam za real-time sistem sa zagarantovanim deadline-om?",
         "EDF / RM, prioritet po roku, preemption, ..."),
    ]
    gap = Inches(0.15)
    n = len(levels)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for lvl, q, a in levels:
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.12), box_h, fill=C_ACCENT)

        # level label
        add_text_box(
            slide, CONTENT_X + Inches(0.25), y + Inches(0.1),
            Inches(2.2), Inches(0.35),
            text=lvl.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        )
        # question
        add_text_box(
            slide, CONTENT_X + Inches(2.6), y + Inches(0.1),
            CONTENT_W - Inches(2.8), Inches(0.55),
            text=f"Q: {q}", size=FS_BODY_SM, bold=True, color=C_PRIMARY,
        )
        # answer
        add_text_box(
            slide, CONTENT_X + Inches(2.6), y + Inches(0.65),
            CONTENT_W - Inches(2.8), box_h - Inches(0.7),
            text=f"A: {a}", size=FS_BODY_SM, italic=True, color=C_TEXT,
        )

        y += box_h + gap


def body_luka_bloom_pokrivenost(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Bloomova pokrivenost (Luka)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Metrika", "strong"),
            "Da li su zastupljeni svi Bloomovi nivoi?",
            "Koliko aktivnosti po nivou?",
            "Da li su ishodi mapirani na više od jednog nivoa?",
            "",
            ("Pristup", "strong"),
            "LLM eksplicitno traži po svakom nivou",
            "Sistem ne prelazi na sledeći nivo dok prethodni nije mapiran",
            "Dedup po (ishod, aktivnost_tip)",
            "",
            ("Nalaz", "strong"),
            "Najbolji rezultati: Pamtiti i Razumeti",
            "Slabiji: Evaluirati i Kreirati - često se 'dotiknu' samo povrsno",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_image_placeholder(
        slide, right_x, CONTENT_Y, half_w, CONTENT_H,
        label="Bar chart - broj aktivnosti po Bloomovom nivou",
    )


def body_stefan_ui_robust(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="UI detekcija - šta sve može da pukne",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Vision model: Qwen 2.5 VL", "strong"),
            "Trazi UI element po opisu (lokacija + box)",
            "",
            ("Ograničenja", "strong"),
            "Non-standard UI frameworks (custom controls)",
            "Brzo menjajuce interfejse (animacije, lazy load)",
            "Nizak kontrast / non-English labele",
            "Multi-monitor setup - samo primarni",
            "",
            ("Posledice", "strong"),
            "Step fail - sistem zaustavlja izvršenje",
            "Video se završava 'kratko' (do tačke prekida)",
            "Logujemo gde je tačno puklo - olaksava debug",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Testirane aplikacije",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    apps = [
        ("Visual Studio 2022", "Stabilno"),
        ("VS Code",            "Stabilno"),
        ("Chrome / Firefox / Edge / Opera", "Stabilno (web)"),
        ("Notepad / Notepad++","Stabilno"),
        ("Custom WinForms",    "Nestabilno"),
        ("Web app sa lazy UI", "Nestabilno"),
    ]
    n = len(apps)
    gap = Inches(0.06)
    apps_area_h = Inches(3.4)
    row_h = (apps_area_h - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for app, status in apps:
        color = C_PRIMARY if "Stabilno" in status else C_ACCENT
        add_filled_rect(slide, right_x, y, half_w, row_h,
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.08), row_h, fill=color)
        add_text_box(
            slide, right_x + Inches(0.2), y, half_w * 0.6, row_h,
            text=app, size=FS_BODY_SM, bold=True, color=C_PRIMARY,
            anchor=MSO_ANCHOR.MIDDLE,
        )
        add_text_box(
            slide, right_x + half_w * 0.62, y, half_w * 0.36, row_h,
            text=status, size=FS_BODY_SM, color=color,
            anchor=MSO_ANCHOR.MIDDLE, align=PP_ALIGN.RIGHT,
        )
        y += row_h + gap

    # Benchmark referenca - OSWorld
    bench_y = y + Inches(0.15)
    bench_h = CONTENT_Y + CONTENT_H - bench_y - Inches(0.1)
    add_filled_rect(slide, right_x, bench_y, half_w, bench_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_filled_rect(slide, right_x, bench_y, Inches(0.1), bench_h, fill=C_ACCENT)
    add_text_box(
        slide, right_x + Inches(0.2), bench_y + Inches(0.06),
        half_w - Inches(0.3), Inches(0.3),
        text="BENCHMARK REFERENCA - OSWorld", size=FS_SECTION_LABEL,
        bold=True, color=C_ACCENT,
    )
    add_text_box(
        slide, right_x + Inches(0.2), bench_y + Inches(0.36),
        half_w - Inches(0.3), bench_h - Inches(0.4),
        text=("Anthropic Computer Use (10/2024): 14.9%   |   "
              "Claude 3.5 Sonnet (12/2024): ~22%   |   "
              "Najbolji open-source vision agent: ~12%.\n"
              "Nas pristup nije autonomni agent već eksplicitan plan - "
              "izbegavamo open-ended traženje sledećeg klika."),
        size=FS_BODY_SM, color=C_TEXT,
    )


def body_pdf_coverage(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="PDF Coverage Tracking (Uroš)",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Šta meri", "strong"),
            "Koje su stranice PDF-a pokrivene pitanjima?",
            "Ponderisana pokrivenost po broju znakova",
            "Sustinska pokrivenost (izuzima skoro prazne stranice)",
            "",
            ("Vizualno - heatmap", "strong"),
            "Bar po stranici, visina = broj znakova",
            "Boja = da li je stranica pokrivena pitanjima",
            "Lista suštinskih stranica bez pitanja",
            "",
            ("Posledica", "strong"),
            "Sistem ti sam kaže gde da generises jos pitanja",
            "Heatmap je 'second pair of eyes' za nastavnika",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_image_placeholder(
        slide, right_x, CONTENT_Y, half_w, CONTENT_H,
        label="Coverage heatmap - screenshot iz UI-ja",
    )


def body_granice_hitl(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Granice automatizacije",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "Definisanje ishoda - zahteva pedagošku procenu",
            "Procena težine pitanja u kontekstu grupe - dinamički kontekst",
            "Procena EA pitanja - često dvosmislena",
            "Recenzija video tutorijala - da li je didaktički dobar",
            "Konačno korigovanje halucinacija",
            "",
            ("Sve ovo radi - čovek", "strong"),
            "Sistem snima pripremu (do 80% rada)",
            "Čovek validira i fino podešava",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Human-in-the-loop u nasem pipeline-u",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    flow = [
        ("Sistem generiše",   "Struktura, pitanja, video"),
        ("Nastavnik pregleda","source_line + ontology anchor + coverage"),
        ("Nastavnik koriguje","Edituje pitanja, raspoređuje aktivnosti"),
        ("Sistem belezi",     "Promene u DB - feedback za buduće generisanje"),
    ]
    gap = Inches(0.12)
    n = len(flow)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for i, (lbl, desc) in enumerate(flow):
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.08), box_h,
                        fill=(C_ACCENT if i % 2 == 0 else C_PRIMARY))

        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.1),
            half_w - Inches(0.3), Inches(0.35),
            text=f"{i+1}. {lbl}", size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.45),
            half_w - Inches(0.3), box_h - Inches(0.5),
            text=desc, size=FS_BODY_SM, color=C_TEXT,
        )
        y += box_h + gap


def body_seg_c_kljucna(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.4), CONTENT_W, Inches(0.4),
        text="Ključna poruka segmenta C",
        size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
        align=PP_ALIGN.CENTER,
    )
    box_w = CONTENT_W * 0.85
    box_x = CONTENT_X + (CONTENT_W - box_w) / 2
    box_y = CONTENT_Y + Inches(1.1)
    box_h = Inches(2.6)
    add_filled_rect(slide, box_x, box_y, box_w, box_h,
                    fill=C_CARD_BG, line=C_LIGHT_LINE)
    add_horizontal_line(slide, box_x, box_y, box_w, color=C_PRIMARY,
                        height=Emu(30000))
    add_text_box(
        slide, box_x + Inches(0.4), box_y + Inches(0.4), box_w - Inches(0.8),
        box_h - Inches(0.8),
        text=("Kvalitet automatski generisanog sadržaja opada sa porastom "
              "kognitivne složenosti - to nije bug, već strukturna posledica, "
              "i upravo zato je human-in-the-loop obavezan."),
        size=22, italic=True, color=C_PRIMARY, align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


# --- SEGMENT D helperi -------------------------------------------------------

def body_sinteza_arhitektura(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Integrisana arhitektura - predlog",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            ("Sloj 1 - Course Planner (Luka)", "strong"),
            "Strukturu kursa generiše iz ishoda + PDF-a",
            "Izlaz: JSON / OWL hijerarhija aktivnosti",
            "",
            ("Sloj 2 - Knowledge Layer (Uroš)", "strong"),
            "ConceptRelationship ontologija iz lekcija",
            "Source-of-truth za relacije pojmova",
            "",
            ("Sloj 3 - Assessment (Uroš)", "strong"),
            "SOLO pitanja, ankored na knowledge layer",
            "",
            ("Sloj 4 - Demonstration (Stefan)", "strong"),
            "Praktiche aktivnosti -> NL instrukcija -> video",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)

    # Vertikalni 4-sloj dijagram sa strelicama
    layers = [
        ("1. Course Planner", "Luka",
         "Ulaz: ishodi + PDF", C_PRIMARY),
        ("2. Knowledge Layer", "Uroš",
         "ConceptRelationship ontologija", C_PRIMARY),
        ("3. Assessment + Demonstration", "Uroš + Stefan",
         "SOLO pitanja  |  OWL-validiran video", C_ACCENT),
        ("4. Human-in-the-loop", "Nastavnik",
         "Recenzija, korekcija, feedback", C_PRIMARY),
    ]
    n = len(layers)
    arrow_h = Inches(0.3)
    avail_h = CONTENT_H - Inches(0.1)
    box_h = (avail_h - (n - 1) * arrow_h) / n
    y = CONTENT_Y

    for i, (title, who, sub, color) in enumerate(layers):
        # Box
        add_filled_rect(slide, right_x, y, half_w, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, right_x, y, Inches(0.12), box_h, fill=color)
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.1),
            half_w * 0.65, Inches(0.4),
            text=title, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, right_x + half_w * 0.7, y + Inches(0.1),
            half_w * 0.28, Inches(0.4),
            text=who.upper(), size=FS_SECTION_LABEL, bold=True, color=C_ACCENT,
            align=PP_ALIGN.RIGHT,
        )
        add_text_box(
            slide, right_x + Inches(0.25), y + Inches(0.5),
            half_w - Inches(0.35), box_h - Inches(0.5),
            text=sub, size=FS_BODY_SM, color=C_TEXT,
        )

        # Strelica dole (osim za poslednji)
        if i < n - 1:
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.DOWN_ARROW,
                right_x + half_w / 2 - Inches(0.15), y + box_h + Inches(0.02),
                Inches(0.3), arrow_h - Inches(0.04),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = C_ACCENT
            arrow.line.fill.background()

        y += box_h + arrow_h


def body_tok_podataka(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="Tok podataka kroz integrisan sistem",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    steps = [
        ("Nastavnik ucitava materijal",
         "PDF lekcije, lista ishoda, opis kursa"),
        ("Course Planner gradi strukturu",
         "Moduli -> koncepti -> ishodi -> aktivnosti (Bloom-aware)"),
        ("Knowledge Layer ekstrahuje ontologiju",
         "Multi-pass LLM ekstrakcija ConceptRelationship-ova"),
        ("Assessment generiše SOLO pitanja",
         "Anchor iz Knowledge Layer-a, PS4 prompt, source_line"),
        ("Demonstration generiše video za praktičnu aktivnost",
         "NL instrukcija -> OWL plan -> validacija -> izvršenje + snimanje"),
        ("Nastavnik pregleda i koriguje",
         "Coverage heatmap, source_line provera, video review"),
    ]
    n = len(steps)
    gap = Inches(0.1)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for i, (h, sub) in enumerate(steps):
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, box_h,
                        fill=C_CARD_BG if i % 2 == 0 else C_WHITE,
                        line=C_LIGHT_LINE)
        # broj
        add_filled_rect(slide, CONTENT_X, y, Inches(0.7), box_h,
                        fill=C_PRIMARY)
        add_text_box(slide, CONTENT_X, y, Inches(0.7), box_h,
                     text=f"{i+1}", size=24, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

        add_text_box(
            slide, CONTENT_X + Inches(0.85), y + Inches(0.05),
            CONTENT_W - Inches(1.0), Inches(0.35),
            text=h, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.85), y + Inches(0.4),
            CONTENT_W - Inches(1.0), box_h - Inches(0.4),
            text=sub, size=FS_BODY_SM, color=C_MUTED,
        )
        y += box_h + gap


def body_veze_drugim_temama(slide):
    cols = [
        ("Deep Reinforcement Learning",
         "Adaptivno biranje pitanja i puta učenja",
         ["State: profil studenta + istorija odgovora",
          "Action: izbor narednog pitanja / aktivnosti",
          "Reward: napredak na SOLO / Bloomovoj skali",
          "Veza sa temom DRL for ITS sa liste predmeta"]),
        ("Recommender Systems",
         "Personalizacija puteva učenja",
         ["Sadržaj-based: ontology anchor sluzi kao feature",
          "Collaborative: studenti sa slicnim odgovorima",
          "Hibridno: koristi obe signala",
          "Veza sa temom Recommenders for ITS"]),
    ]
    gap = Inches(0.3)
    col_w = (CONTENT_W - gap) / 2

    for i, (title, sub, points) in enumerate(cols):
        x = CONTENT_X + i * (col_w + gap)
        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.1),
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=C_ACCENT)

        pad = Inches(0.25)
        cur_y = CONTENT_Y + Inches(0.3)
        add_text_box(slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.5),
                     text=title, size=20, bold=True, color=C_PRIMARY)
        cur_y += Inches(0.55)
        add_text_box(slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.4),
                     text=sub, size=FS_BODY_SM, italic=True, color=C_ACCENT)
        cur_y += Inches(0.4)
        add_horizontal_line(slide, x + pad, cur_y, Inches(0.6),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.2)
        add_bullet_list(slide, x + pad, cur_y, col_w - 2 * pad,
                        CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.3),
                        points, size=FS_BODY_SM)


def body_otvorena_pitanja(slide):
    pitanja = [
        ("Kvalitet distraktora na EA nivou",
         "Kako podići uverljivost bez gubitka SOLO 'tipa' greske?"),
        ("Real-world evaluacija",
         "Sistem nije jos testiran u stvarnoj učionici - to je sledeći korak."),
        ("Multimodalna procena znanja",
         "Kako uracunati ne samo tekstualne odgovore već i praktične demonstracije?"),
        ("Stabilnost vision modela",
         "Custom UI ostaju izazov; integracija sa accessibility API-jima?"),
        ("Skaliranje na više predmeta",
         "Da li se ontologija prenosi između domena ili je svaki predmet svoja TBox?"),
    ]
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, CONTENT_W, Inches(0.4),
        text="Otvorena pitanja", size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )

    n = len(pitanja)
    gap = Inches(0.12)
    box_h = (CONTENT_H - Inches(0.5) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.5)

    for i, (h, sub) in enumerate(pitanja):
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.5), box_h, fill=C_ACCENT)
        add_text_box(slide, CONTENT_X, y, Inches(0.5), box_h,
                     text=f"Q{i+1}", size=16, bold=True, color=C_WHITE,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(
            slide, CONTENT_X + Inches(0.65), y + Inches(0.05),
            CONTENT_W - Inches(0.8), Inches(0.4),
            text=h, size=FS_BODY, bold=True, color=C_PRIMARY,
        )
        add_text_box(
            slide, CONTENT_X + Inches(0.65), y + Inches(0.4),
            CONTENT_W - Inches(0.8), box_h - Inches(0.4),
            text=sub, size=FS_BODY_SM, color=C_TEXT,
        )
        y += box_h + gap


def body_etika_ai_act(slide):
    """Etika, privatnost, EU AI Act - kako nas dizajn već adresira high-risk zahteve."""
    cards = [
        ("Privatnost",
         "Lokalni LLM (Ollama) nije slučajan izbor",
         [
             "PDF nastavnog materijala često sadrzi IP fakulteta",
             "Studentski podaci ne napuštaju mašinu",
             "Zero per-call cost - skalabilno za institucije",
         ],
         C_PRIMARY),
        ("Transparentnost",
         "source_line + ontology anchor = audit trail",
         [
             "Svako pitanje ima izvorni navod iz materijala",
             "Svaka relacija ima tipiziranu vezu u ontologiji",
             "Recenzent može da prati 'zašto' iza svake odluke",
         ],
         C_ACCENT),
        ("EU AI Act (2024/2025)",
         "Obrazovni AI je klasifikovan kao high-risk",
         [
             "Striktna transparentnost: dokumentovati izvore i logiku",
             "Obavezan ljudski oversight (human-in-the-loop)",
             "Audit trail - ono sto nas dizajn već ima 'iz kutije'",
         ],
         C_PRIMARY),
    ]

    gap = Inches(0.25)
    col_w = (CONTENT_W - 2 * gap) / 3

    for i, (title, sub, points, color) in enumerate(cards):
        x = CONTENT_X + i * (col_w + gap)
        add_filled_rect(slide, x, CONTENT_Y, col_w, CONTENT_H - Inches(0.5),
                        fill=C_WHITE, line=C_LIGHT_LINE)
        add_filled_rect(slide, x, CONTENT_Y, col_w, Inches(0.08), fill=color)

        pad = Inches(0.22)
        cur_y = CONTENT_Y + Inches(0.25)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.4),
            text=title, size=20, bold=True, color=C_PRIMARY,
        )
        cur_y += Inches(0.5)
        add_text_box(
            slide, x + pad, cur_y, col_w - 2 * pad, Inches(0.7),
            text=sub, size=FS_BODY_SM, italic=True, color=color,
        )
        cur_y += Inches(0.7)
        add_horizontal_line(slide, x + pad, cur_y, Inches(0.5),
                            color=C_ACCENT, height=Emu(15000))
        cur_y += Inches(0.18)
        add_bullet_list(
            slide, x + pad, cur_y, col_w - 2 * pad,
            CONTENT_H - (cur_y - CONTENT_Y) - Inches(0.5),
            points, size=FS_BODY_SM,
        )

    # Donji red - sažeta poruka
    y2 = CONTENT_Y + CONTENT_H - Inches(0.4)
    add_text_box(
        slide, CONTENT_X, y2, CONTENT_W, Inches(0.35),
        text="Nas dizajn nije optimizovan za AI Act - ali ga zadovoljava prirodno, "
             "zbog izbora ontology grounding-a i lokalnih modela.",
        size=FS_BODY_SM, italic=True, color=C_MUTED, align=PP_ALIGN.CENTER,
    )


def body_buduci_rad(slide):
    half_w = CONTENT_W / 2 - Inches(0.2)
    add_text_box(
        slide, CONTENT_X, CONTENT_Y, half_w, Inches(0.4),
        text="Naredni koraci - kratki rok",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, CONTENT_X, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "Spojiti tri repozitorijuma u jedan monorepo",
            "Definisati API kontrakte između slojeva",
            "Zajednička ontološka osnova (deljena seed_ontology)",
            "Demo end-to-end pipe za jedan predmet",
            "",
            ("Konferencija Sinteza 2026", "strong"),
            "Jedan rad iz tima već publikovan",
            "Cilj: drugi rad o integraciji",
        ],
        size=FS_BODY_SM,
    )

    right_x = CONTENT_X + half_w + Inches(0.4)
    add_text_box(
        slide, right_x, CONTENT_Y, half_w, Inches(0.4),
        text="Naredni koraci - duži rok",
        size=FS_SUBTITLE, bold=True, color=C_PRIMARY,
    )
    add_bullet_list(
        slide, right_x, CONTENT_Y + Inches(0.5), half_w, CONTENT_H - Inches(0.6),
        [
            "Pilot u stvarnoj učionici",
            "Adaptivno biranje pitanja (DRL prototip)",
            "Personalizacija puta učenja (Recommender)",
            "Multimodalna procena znanja (audio / kod / video)",
            "Skaliranje na više predmeta",
            "",
            ("Etika i privatnost", "strong"),
            "Lokalni LLM kao osnovna opcija",
            "Pravo nastavnika na konačnu reč",
        ],
        size=FS_BODY_SM,
    )


def body_zakljucak(slide):
    add_text_box(
        slide, CONTENT_X, CONTENT_Y + Inches(0.2), CONTENT_W, Inches(0.4),
        text="Šta smo pokazali", size=FS_SECTION_LABEL, bold=True,
        color=C_ACCENT, align=PP_ALIGN.CENTER,
    )

    points = [
        ("LLM sam nije dovoljan", "halucinacije + nedostatak strukture"),
        ("Pedagoška taksonomija kalibrise", "Bloom + SOLO kao kompas"),
        ("Ontologija filtrira", "OWL/SPARQL + source_line + ontology anchor"),
        ("Human-in-the-loop ostaje obavezan", "naročito na višim SOLO nivoima"),
    ]
    gap = Inches(0.18)
    n = len(points)
    box_h = (CONTENT_H - Inches(1.2) - (n - 1) * gap) / n
    y = CONTENT_Y + Inches(0.7)

    for i, (h, sub) in enumerate(points):
        add_filled_rect(slide, CONTENT_X, y, CONTENT_W, box_h,
                        fill=C_CARD_BG, line=C_LIGHT_LINE)
        add_filled_rect(slide, CONTENT_X, y, Inches(0.1), box_h, fill=C_ACCENT)
        add_text_box(slide, CONTENT_X + Inches(0.25), y + Inches(0.1),
                     CONTENT_W - Inches(0.4), Inches(0.4),
                     text=h, size=FS_BODY, bold=True, color=C_PRIMARY)
        add_text_box(slide, CONTENT_X + Inches(0.25), y + Inches(0.45),
                     CONTENT_W - Inches(0.4), box_h - Inches(0.5),
                     text=sub, size=FS_BODY_SM, color=C_TEXT)
        y += box_h + gap


def body_literatura_lista(slide, refs):
    """Lista referenci (po kategorijama)."""
    cur_y = CONTENT_Y
    avail_h = CONTENT_H

    for cat, items in refs:
        add_text_box(
            slide, CONTENT_X, cur_y, CONTENT_W, Inches(0.32),
            text=cat, size=FS_BODY, bold=True, color=C_ACCENT,
        )
        cur_y += Inches(0.35)
        for ref in items:
            add_text_box(
                slide, CONTENT_X + Inches(0.2), cur_y,
                CONTENT_W - Inches(0.2), Inches(0.45),
                text="- " + ref, size=FS_BODY_SM, color=C_TEXT,
            )
            cur_y += Inches(0.48)
        cur_y += Inches(0.1)


# =============================================================================
# IZGRADNJA PREZENTACIJE
# =============================================================================

def build_presentation(out_path):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Pratimo broj slidova za footer
    slide_specs = []

    # ---- Naslovna ----
    slide_specs.append(("title", {}))

    # ---- Uvod (predmet + tim + agenda) ----
    slide_specs.append(("content", dict(
        page_label=None,
        title="Tim",
        subtitle="Tri projekta sa predmeta SOTIS - jedan integrisan ITS pipeline",
        body_func=body_team,
    )))
    slide_specs.append(("content", dict(
        page_label=None,
        title="Predmet i kontekst",
        subtitle="Napredne tehnike računarske inteligencije",
        body_func=body_predmet_context,
    )))
    slide_specs.append(("content", dict(
        page_label=None,
        title="Agenda",
        subtitle="5 segmenata - tematska organizacija, ne po projektima",
        body_func=body_agenda,
    )))

    # ============ SEGMENT 0 - UVOD I MOTIVACIJA ============
    slide_specs.append(("divider", dict(
        segment_label="Segment 0",
        title="Uvod i motivacija",
        subtitle="Polazimo od trenutnog stanja - šta nastavnik radi i zašto je to skupo.",
        time_str="10 min  /  4 slajda",
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Uvod",
        title="Problem - šta nastavnik radi ručno",
        subtitle="Vreme i kognitivni napor su jaka skala bola.",
        body_func=body_problem_stats,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Uvod",
        title="LLM otvara vrata, ali ne sam",
        subtitle="Šta dobijemo - i šta sve uneso skripa.",
        body_func=body_llm_kao_odgovor,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Zašto sad",
        title="Zašto bas sad? - tri empirijska signala",
        subtitle="Od Bloomovog '2 Sigma' problema do mainstream LLM tutora.",
        body_func=body_why_now,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Landscape",
        title="ITS landscape - gde se mi uklapamo",
        subtitle="Postojeci sistemi serviraju materijal; mi automatizujemo njegovo generisanje.",
        body_func=body_its_landscape,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Vizija",
        title="Vizija - ITS pipeline od ishoda do video tutorijala",
        subtitle="5 koraka, end-to-end, bez rucne intervencije između.",
        body_func=body_vizija_pipeline,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Komponente",
        title="Tri komponente jednog pipeline-a",
        subtitle="Svaki član tima radi jedan deo - zajedno čine celinu.",
        body_func=body_tri_komponente,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment 0 - Ključna poruka",
        title="Ključna poruka prezentacije",
        body_func=body_kljucna_poruka_uvod,
    )))

    # ============ SEGMENT A - LLM PROMPTING ============
    slide_specs.append(("divider", dict(
        segment_label="Segment A",
        title="LLM prompting i pedagoške taksonomije",
        subtitle="Hijerarhijska dekompozicija + formalna taksonomija u promptu = pedagoški kalibrisan LLM.",
        time_str="25 min  /  11 slajdova",
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Bloom",
        title="Bloomova taksonomija - 6 nivoa misaonih operacija",
        subtitle="Glagol akcije: pamtiti, razumeti, primeniti, analizirati, evaluirati, kreirati.",
        body_func=body_bloom_pregled,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - SOLO",
        title="SOLO taksonomija - strukturalna složenost odgovora",
        subtitle="Od jedne činjenice do generalizacije i transfera znanja.",
        body_func=body_solo_pregled,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Poređenje",
        title="Bloom vs SOLO - različiti pogledi, komplementarni",
        body_func=body_bloom_vs_solo,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Luka",
        title="Luka - hijerarhijska struktura kursa",
        subtitle="LLM razlaže kurs u stablo modula, koncepata, ishoda i aktivnosti.",
        body_func=body_luka_struktura,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Luka",
        title="Luka - prompt template + state-of-generation",
        body_func=body_luka_prompt_template,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Uroš",
        title="Uroš - SOLO Quiz Generator i PS4 prompt template",
        body_func=body_uros_solo_ps4,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Uroš",
        title="Two-pass generisanje za Extended Abstract",
        subtitle="Razdvajamo 'šta pitati' od 'kako napraviti dobre distraktore'.",
        body_func=body_two_pass_ea,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Uroš",
        title="Source-line citation - anti-halucinacija",
        body_func=body_source_line,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Stefan",
        title="Stefan - hijerarhijska dekompozicija zadataka",
        subtitle="Task -> Step -> Action: tri nivoa razgradnje prirodno-jezičke instrukcije.",
        body_func=body_stefan_dekompozicija,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Tradeoff",
        title="Prompt engineering vs Fine-tuning",
        body_func=body_prompt_vs_finetune,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Tradeoff",
        title="Lokalni vs cloud LLM-ovi - mi koristimo oba",
        body_func=body_local_vs_cloud,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment A - Ključna poruka",
        title="Ključna poruka segmenta A",
        body_func=body_seg_a_kljucna,
    )))

    # ============ SEGMENT B - ONTOLOGIJE ============
    slide_specs.append(("divider", dict(
        segment_label="Segment B",
        title="Ontologije i formalni grounding",
        subtitle="Ontologija je ugovor između LLM-a i sistema - sve sto LLM proizvede mora se mapirati.",
        time_str="25 min  /  10 slajdova",
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Osnove",
        title="OWL - Web Ontology Language",
        subtitle="Klase, instance, relacije, ograničenja, mas i mas serijalizacija.",
        body_func=body_owl_uvod,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Osnove",
        title="RDF i SPARQL - tripleti i upiti",
        subtitle="Sve činjenice su tripleti; sve upite pišemo nad gradom triplet-a.",
        body_func=body_rdf_sparql,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Luka",
        title="Luka - JSON serijalizacija strukture kursa",
        subtitle="LLM lakše generiše validan JSON nego direktan OWL - JSON -> OWL je determinizam.",
        body_func=body_luka_json_primer,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Uroš",
        title="Uroš - ConceptRelationship ontologija",
        subtitle="Relacije između learning objekata - osnova za relacionalna i EA pitanja.",
        body_func=body_uros_ontologija,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Uroš",
        title="Uroš - SPARQL upit za 'ontology anchor'",
        subtitle="Pitanje se ne pravi 'iz vazduha' - bira se konkretna veza i ona postaje kontekst.",
        body_func=body_ontology_anchor,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Stefan",
        title="Stefan - OWL ontologija izvršnog plana",
        subtitle="Plan koji LLM generiše mora se mapirati na klase Task / Step / Action.",
        body_func=body_stefan_owl,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Stefan",
        title="Stefan - pravila validacije plana (SPARQL)",
        subtitle="Pre nego sto se ista izvrši, SPARQL upiti proveravaju strukturna pravila.",
        body_func=body_validacija_pravila,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Anti-halucinacija",
        title="Anti-halucinacioni mehanizmi - 3 ugla",
        subtitle="Svaki član tima pokriva drugačiju vrstu greske LLM-a.",
        body_func=body_anti_halucinacije,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Anti-halucinacija",
        title="Anatomija halucinacije - jedan konkretan primer",
        subtitle="Levo: bez check-a, halucinacija prolazi. Desno: source_line + fuzzy match je hvataju.",
        body_func=body_anatomija_halucinacije,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment B - Ključna poruka",
        title="Ključna poruka segmenta B",
        body_func=body_seg_b_kljucna,
    )))

    # ============ SEGMENT C - EVALUACIJA ============
    slide_specs.append(("divider", dict(
        segment_label="Segment C",
        title="Evaluacija, kvalitet i ograničenja",
        subtitle="Kvalitet opada sa kognitivnom složenošću - to opravdava human-in-the-loop.",
        time_str="25 min  /  9 slajdova",
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - A-priori validacija",
        title="Pre nego sto student vidi pitanje - 3 sloja provere",
        subtitle="Cela 'baterija' tehnika iz literature: prevent + generate + detect, sve pre čoveka.",
        body_func=body_validation_stack,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Haladyna",
        title="Haladyna item-writing pravila - prevent + detect hybrid",
        subtitle="Ista pravila u promptu i u lint-u - lint flag se može kvotirati nazad u prompt rule.",
        body_func=body_haladyna_lint,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - LLM-as-judge",
        title="LLM kao kritičar samog sebe - CoVe, Solvability, SOLO judge, IOC",
        subtitle="Posle generisanja, pre čoveka: drugi LLM (ili isti sa drugacijim promptom) procenjuje izlaz.",
        body_func=body_llm_kao_kriticar,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Pilot",
        title="Pilot evaluacija SOLO pitanja - ljudska anotacija",
        subtitle="Posle automatske validacije: tim + nezavisan recenzent, 4 dimenzije, Likert 1-5.",
        body_func=body_pilot_metodologija,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Rezultati",
        title="Rezultati - kvalitet po SOLO nivou",
        subtitle="Validnost i kvalitet distraktora padaju kako se ide ka EA nivou.",
        body_func=body_kvalitet_po_nivou,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Primeri",
        title="Primeri SOLO pitanja iz domena Operativnih sistema",
        body_func=body_primer_solo_pitanja,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Luka",
        title="Luka - Bloomova pokrivenost generisanih kurseva",
        subtitle="Vidimo gde sistem dobro pokriva nivoe, a gde tek 'dotiče' viša mišljenja.",
        body_func=body_luka_bloom_pokrivenost,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Stefan",
        title="Stefan - robustnost izvršenja video tutorijala",
        subtitle="Vision modeli su odlični, ali ne savršeni - non-standard UI je glavna bolna tačka.",
        body_func=body_stefan_ui_robust,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Uroš",
        title="PDF Coverage Tracking - heatmap pokrivenosti",
        subtitle="Sistem sam ti kaže gde da generises jos pitanja.",
        body_func=body_pdf_coverage,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - HITL",
        title="Granice automatizacije i human-in-the-loop",
        subtitle="Sistem snima 80% rada; čovek validira i fino podešava.",
        body_func=body_granice_hitl,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment C - Ključna poruka",
        title="Ključna poruka segmenta C",
        body_func=body_seg_c_kljucna,
    )))

    # ============ SEGMENT D - SINTEZA ============
    slide_specs.append(("divider", dict(
        segment_label="Segment D",
        title="Sinteza, budući rad, zaključak",
        subtitle="Kako tri komponente postaju jedan sistem - i šta sledi.",
        time_str="10-15 min  /  7 slajdova",
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Arhitektura",
        title="Predlog integrisane arhitekture - 4 sloja",
        subtitle="Sloj po sloj: planiranje -> znanje -> procena -> demonstracija.",
        body_func=body_sinteza_arhitektura,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Arhitektura",
        title="Tok podataka kroz integrisan sistem",
        subtitle="Sest koraka, od ucitavanja materijala do recenzije nastavnika.",
        body_func=body_tok_podataka,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Veze",
        title="Veze sa drugim temama predmeta",
        subtitle="DRL i Recommender Systems prirodno produžavaju ovaj pipeline.",
        body_func=body_veze_drugim_temama,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Buduci rad",
        title="Otvorena pitanja",
        body_func=body_otvorena_pitanja,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Etika",
        title="Etika, privatnost i EU AI Act",
        subtitle="Obrazovni AI je high-risk po EU AI Act-u - nas dizajn već adresira kljucne zahteve.",
        body_func=body_etika_ai_act,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Buduci rad",
        title="Naredni koraci - kratki i duži rok",
        body_func=body_buduci_rad,
    )))
    slide_specs.append(("content", dict(
        page_label="Segment D - Zakljucak",
        title="Zakljucak",
        subtitle="Vracamo se na početnu poruku.",
        body_func=body_zakljucak,
    )))
    slide_specs.append(("quote", dict(
        page_label="Zaključna misao",
        quote=("LLM + ontologija + pedagoška taksonomija = "
               "pouzdano, pedagoški kalibrisano ITS okruženje koje je "
               "više od zbira svojih delova."),
        attribution="ključna poruka prezentacije",
    )))

    # ---- Hvala ----
    slide_specs.append(("thanks", {}))

    # ---- Literatura ----
    literatura_1 = [
        ("Pedagoške taksonomije", [
            "Bloom, B.S. (1984). The 2 Sigma Problem. Educational Researcher, 13(6), 4-16.",
            "Biggs, J.B. & Collis, K.F. (1982). Evaluating the Quality of Learning: The SOLO Taxonomy. Academic Press.",
            "Anderson, L.W. & Krathwohl, D.R. (2001). A Taxonomy for Learning, Teaching, and Assessing. Longman.",
            "Lister, R., et al. (2006). Not Seeing the Forest for the Trees: Novice Programmers and the SOLO Taxonomy. ACM SIGCSE Bulletin, 38(3).",
        ]),
        ("LLM u obrazovanju", [
            "Kasneci, E., et al. (2023). ChatGPT for good? Learning and Individual Differences, 103.",
            "Yan, L., et al. (2024). Practical and ethical challenges of LLMs in education. BJET, 55(1).",
            "Wang, S., et al. (2024). Large Language Models for Education: A Survey. arXiv:2403.18105.",
            "Kestin, G., et al. (2024). AI Tutoring Outperforms In-Class Active Learning. Preprint, PhysPort.",
            "Mollick, E. & Mollick, L. (2023). Assigning AI: Seven Approaches for Students. Wharton, SSRN 4475995.",
        ]),
    ]
    literatura_2 = [
        ("Prompting i generisanje pitanja", [
            "Scaria, N., et al. (2024). PS4 - Automated Educational QG at Bloom's Skill Levels. arXiv:2408.04394.",
            "Wei, J., et al. (2022). Chain-of-Thought Prompting Elicits Reasoning in LLMs. NeurIPS 2022. arXiv:2201.11903.",
            "Sweller, J. & Cooper, G.A. (1985). The use of worked examples as a substitute for problem solving. Cognition & Instruction 2(1).",
            "Wang, X., et al. (2023). Self-Consistency Improves CoT Reasoning. ICLR 2023. arXiv:2203.11171.",
            "Kurdi, G., et al. (2020). A Systematic Review of Automatic Question Generation. IJAIED, 30(1).",
        ]),
        ("Distractor generation i misconception theory", [
            "Bitew, S.K., et al. (2023). Distractor Generation with Predictive Prompting and LLMs. arXiv:2307.16338.",
            "Sadler, P.M. (1998). Psychometric Models of Student Conceptions in Science. JRST, 35(3).",
            "Liang, C., et al. (2018). Distractor Generation for MCQ Using Learning to Rank. NAACL Workshop.",
            "Aldabe, I., et al. (2009). Automatic Distractor Generation for Domain Specific Texts. LNCS 6233.",
        ]),
    ]
    literatura_3 = [
        ("Anti-halucinacija i RAG", [
            "Lewis, P., et al. (2020). Retrieval-Augmented Generation for Knowledge-Intensive NLP. NeurIPS 2020. arXiv:2005.11401.",
            "Dhuliawala, S., et al. (2023). Chain-of-Verification (CoVe). ACL Findings 2024. arXiv:2309.11495.",
            "Pal, A., et al. (2023). Med-HALT: Medical Domain Hallucination Test for LLMs. CoNLL 2023.",
            "Zhang, T., et al. (2020). BERTScore: Evaluating Text Generation with BERT. ICLR 2020. arXiv:1904.09675.",
            "Zheng, L., et al. (2023). Judging LLM-as-a-Judge with MT-Bench. NeurIPS 2023. arXiv:2306.05685.",
        ]),
        ("Item analysis i psihometrija", [
            "Haladyna, T.M., Downing, S.M. & Rodriguez, M.C. (2002). MCQ Item-Writing Guidelines. Applied Measurement in Education, 15(3).",
            "Rovinelli, R.J. & Hambleton, R.K. (1977). Item-Objective Congruence (IOC). Dutch Journal of Educational Research, 2.",
            "Crocker, L. & Algina, J. (1986). Introduction to Classical and Modern Test Theory. Holt, Rinehart and Winston.",
            "Landis, J.R. & Koch, G.G. (1977). Observer Agreement for Categorical Data (Cohen's kappa). Biometrics, 33(1).",
            "Tarrant, M., et al. (2009). Item writing flaws in MCQs in high-stakes nursing assessments. Nurse Education in Practice, 9(3).",
        ]),
    ]
    literatura_4 = [
        ("Ontologije i Semantic Web", [
            "Lin, C.-Y., et al. (2025). KAQG: Knowledge-Graph-Enhanced RAG for Difficulty-Controlled QG. arXiv:2505.07618.",
            "Vesin, B., et al. (2013). Ontology-based architecture with recommendation strategy in Java tutoring (ProTuS). CSIS, 10(1).",
            "Hogan, A., et al. (2021). Knowledge Graphs. ACM Computing Surveys, 54(4).",
        ]),
        ("Intelligent Tutoring Systems", [
            "Mousavinasab, E., et al. (2021). ITS: systematic review. Interactive Learning Environments, 29(1).",
            "Lin, C.-C., et al. (2023). AI in ITS toward sustainable education. Smart Learning Environments, 10(1).",
            "VanLehn, K. (2011). The Relative Effectiveness of Human Tutoring, ITS, and Other Tutoring Systems. Educational Psychologist, 46(4).",
        ]),
        ("Multimodalni agenti i Computer Use", [
            "Anthropic (Oct 2024). Computer use model evaluation on OSWorld: 14.9% baseline.",
            "Xie, T., et al. (2024). OSWorld: Benchmarking Multimodal Agents for Open-Ended Tasks. NeurIPS 2024.",
        ]),
    ]

    slide_specs.append(("references", dict(
        page_label="Literatura - 1 / 4",
        title="Literatura",
        subtitle="Pedagoške taksonomije i LLM u obrazovanju",
        refs=literatura_1,
    )))
    slide_specs.append(("references", dict(
        page_label="Literatura - 2 / 4",
        title="Literatura (nastavak)",
        subtitle="Prompting strategije, QG i misconception theory",
        refs=literatura_2,
    )))
    slide_specs.append(("references", dict(
        page_label="Literatura - 3 / 4",
        title="Literatura (nastavak)",
        subtitle="Anti-halucinacija, RAG, item analysis i psihometrija",
        refs=literatura_3,
    )))
    slide_specs.append(("references", dict(
        page_label="Literatura - 4 / 4",
        title="Literatura (nastavak)",
        subtitle="Ontologije, ITS i multimodalni agenti",
        refs=literatura_4,
    )))

    total_pages = len(slide_specs)

    # ---- Renderovanje ----
    for i, (kind, args) in enumerate(slide_specs, start=1):
        if kind == "title":
            make_title_slide(
                prs,
                title=("Od ishoda učenja do automatizovanog ITS-a"),
                subtitle=("LLM-driven pipeline za generisanje strukture kursa, "
                          "video demonstracija i pitanja zasnovanih na SOLO taksonomiji"),
                team="Uroš Petrašković - Luka Šarić - Stefan Lazarević",
                course="Napredne tehnike računarske inteligencije - Master, FTN UNS",
                date_str="2026.",
            )
        elif kind == "divider":
            make_section_divider(prs, **args)
        elif kind == "content":
            make_content_slide(prs, page_no=i, total_pages=total_pages, **args)
        elif kind == "quote":
            make_quote_slide(prs, page_no=i, total_pages=total_pages, **args)
        elif kind == "thanks":
            make_thank_you_slide(
                prs,
                team="Uroš Petrašković - Luka Šarić - Stefan Lazarević",
            )
        elif kind == "references":
            slide = new_slide(prs)
            add_header(slide,
                       title=args["title"], subtitle=args.get("subtitle"),
                       page_label=args.get("page_label"))
            body_literatura_lista(slide, args["refs"])
            add_footer(slide, page_number=i, total=total_pages)

    prs.save(out_path)
    print(f"Generisano: {out_path}  ({total_pages} slidova)")


if __name__ == "__main__":
    import os
    here = os.path.dirname(os.path.abspath(__file__))
    out = os.path.join(here, "Prezentacija_ITS_pipeline.pptx")
    build_presentation(out)
